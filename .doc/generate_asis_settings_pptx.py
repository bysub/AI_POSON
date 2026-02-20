# -*- coding: utf-8 -*-
"""ASIS VB6 설정 관리 - INI vs DB 분류 전체 목록 PPTX 생성 (v2 - 한글 완전판)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1E, 0x29, 0x3B)
BG_ROW = RGBColor(0x16, 0x20, 0x33)
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)
ACCENT_GREEN = RGBColor(0x4A, 0xDE, 0x80)
ACCENT_ORANGE = RGBColor(0xFB, 0xBF, 0x24)
ACCENT_RED = RGBColor(0xF8, 0x71, 0x71)
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
ACCENT_TEAL = RGBColor(0x2D, 0xD4, 0xBF)
ACCENT_PINK = RGBColor(0xF4, 0x72, 0xB6)
ACCENT_CYAN = RGBColor(0x22, 0xD3, 0xEE)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
TEXT_LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
BORDER_COLOR = RGBColor(0x33, 0x41, 0x55)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=14, color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.font.name = "Malgun Gothic"
    return txBox


def add_badge(slide, left, top, text, color=ACCENT_BLUE, width=Inches(1.2)):
    shape = add_shape(slide, left, top, width, Inches(0.32), fill_color=color)
    shape.text_frame.paragraphs[0].text = str(text)
    shape.text_frame.paragraphs[0].font.size = Pt(10)
    shape.text_frame.paragraphs[0].font.color.rgb = TEXT_WHITE
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "Malgun Gothic"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


def add_stat_card(slide, x, y, number, label, color, w=Inches(1.9), h=Inches(1.3)):
    card = add_shape(slide, x, y, w, h, fill_color=BG_CARD, border_color=color)
    add_text(slide, x + Inches(0.1), y + Inches(0.15), w - Inches(0.2), Inches(0.6),
             str(number), 32, color, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), y + Inches(0.8), w - Inches(0.2), Inches(0.4),
             label, 11, TEXT_GRAY, False, PP_ALIGN.CENTER)
    return card


def add_item_row(slide, x, y, key, desc, default="", key_w=Inches(2.2), desc_w=Inches(3.0), def_w=Inches(1.0), size=9):
    add_text(slide, x, y, key_w, Inches(0.22), key, size, ACCENT_CYAN, False)
    add_text(slide, x + key_w, y, desc_w, Inches(0.22), desc, size, TEXT_LIGHT, False)
    if default:
        add_text(slide, x + key_w + desc_w, y, def_w, Inches(0.22), default, size - 1, TEXT_GRAY, False)


def add_section_header(slide, x, y, title, count, color=ACCENT_BLUE):
    add_text(slide, x, y, Inches(4), Inches(0.35), title, 13, color, True)
    add_badge(slide, x + Inches(4.1), y + Inches(0.02), str(count) + "keys", color, Inches(0.9))


# ================================================================
# SLIDE 1: Title
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_ORANGE)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "POSON POS - ASIS VB6 Legacy", 20, TEXT_GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
         "INI vs DB 설정 분류 전체 목록", 44, TEXT_WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.6),
         "pos_config.ini  |  config.ini  |  DB: POS_Set[101]  |  DB: Office_User", 16, ACCENT_ORANGE, False, PP_ALIGN.CENTER)

cards_data = [
    ("INI (POS 런타임)", "~446", ACCENT_BLUE),
    ("INI (관리프로그램)", "~350", ACCENT_PURPLE),
    ("DB POS_Set", "~78", ACCENT_GREEN),
    ("DB Office_User", "~24", ACCENT_ORANGE),
]
card_w = Inches(2.3)
start_x = Inches(1.7)
for i, (label, value, color) in enumerate(cards_data):
    cx = start_x + i * Inches(2.6)
    add_stat_card(slide, cx, Inches(4.5), value, label, color, card_w, Inches(1.2))

add_text(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
         "POS 런타임 합계: 548항목  |  전체(관리프로그램 포함): ~900항목", 14, TEXT_GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(6.7), Inches(11), Inches(0.4),
         "2026-02-19  |  VB6 ASIS 설정 분석", 11, TEXT_GRAY, False, PP_ALIGN.CENTER)


# ================================================================
# SLIDE 2: Overview - 4대 저장소
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_BLUE)
add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
         "ASIS 설정 저장소 4대 개요", 24, TEXT_WHITE, True)
add_text(slide, Inches(0.5), Inches(0.7), Inches(12), Inches(0.4),
         "VB6 Type: C_Config(300+), Terminal(50+), S_Config(130+), Shop(40+), VAN(20+)", 11, TEXT_GRAY)

# Source 1: pos_config.ini
add_shape(slide, Inches(0.3), Inches(1.25), Inches(6.2), Inches(2.7), border_color=ACCENT_BLUE)
add_badge(slide, Inches(0.5), Inches(1.35), "INI FILE", ACCENT_BLUE, Inches(1.0))
add_text(slide, Inches(1.6), Inches(1.35), Inches(4), Inches(0.3), "pos_config.ini", 14, TEXT_WHITE, True)
add_text(slide, Inches(0.5), Inches(1.75), Inches(5.8), Inches(0.3),
         "POS 런타임 설정 / 로컬파일 / 75섹션, 516+키", 10, TEXT_GRAY)
items1 = [
    ("[Application]", "애플리케이션 정보 (5키)"),
    ("[Sale]", "판매 운영 설정 (10키)"),
    ("[Receipt]", "영수증 정보 (9키)"),
    ("[DataBase]", "DB 연결 정보 (6키)"),
    ("[Terminal]", "터미널 HW 설정 (42키)"),
    ("[Other]", "POS 업무 설정 (75키)"),
    ("[Card]+VAN 13개사", "VAN 결제 설정 (65+키)"),
    ("[Self]", "셀프키오스크 전용 (73키)"),
    ("[FaceCam/JaPan/Selfimg/Self21]", "기타 기기 전용 (20키)"),
    ("[Dual_Msg/POSON2/Color/Font]", "UI/색상/폰트 (32키)"),
    ("[Grid_*] 20+섹션", "그리드 컬럼폭 VB6전용 (~100키)"),
]
for i, (sec, desc) in enumerate(items1):
    ry = Inches(2.05) + Inches(i * 0.17)
    add_text(slide, Inches(0.5), ry, Inches(3.0), Inches(0.18), sec, 8, ACCENT_CYAN, False)
    add_text(slide, Inches(3.5), ry, Inches(2.8), Inches(0.18), desc, 8, TEXT_LIGHT, False)

# Source 2: config.ini
add_shape(slide, Inches(6.8), Inches(1.25), Inches(6.2), Inches(2.7), border_color=ACCENT_PURPLE)
add_badge(slide, Inches(7.0), Inches(1.35), "INI FILE", ACCENT_PURPLE, Inches(1.0))
add_text(slide, Inches(8.1), Inches(1.35), Inches(4), Inches(0.3), "config.ini", 14, TEXT_WHITE, True)
add_text(slide, Inches(7.0), Inches(1.75), Inches(5.8), Inches(0.3),
         "PosManager 전용 / 관리프로그램 / 85+섹션, 1317줄", 10, TEXT_GRAY)
items2 = [
    ("[Application]", "관리프로그램 기본설정 (76키)"),
    ("[Server]", "DB 서버 + 원격서버 (15키)"),
    ("[Version_Info]", "버전 정보 (7키)"),
    ("[BarCodePrint]", "바코드 프린터 공통 (43키)"),
    ("[TTP/SRP/LK-B/LK-P/G5]", "프린터 모델별 좌표 (200+키)"),
    ("[HT_Trans]", "데이터 전송 (7키)"),
    ("[WEBSMS_Sett]", "웹 SMS 설정 (8키)"),
    ("[Customer]", "고객표시기 CDP 포트 (4키)"),
    ("[Length]", "바코드/저울 길이 (2키)"),
    ("[SuSu]", "수수료율 설정 (5키)"),
    ("[기타 10+섹션]", "화면별 옵션 (50+키)"),
]
for i, (sec, desc) in enumerate(items2):
    ry = Inches(2.05) + Inches(i * 0.17)
    add_text(slide, Inches(7.0), ry, Inches(3.0), Inches(0.18), sec, 8, ACCENT_PURPLE, False)
    add_text(slide, Inches(10.0), ry, Inches(2.8), Inches(0.18), desc, 8, TEXT_LIGHT, False)

# Source 3: DB POS_Set[101]
add_shape(slide, Inches(0.3), Inches(4.2), Inches(6.2), Inches(2.9), border_color=ACCENT_GREEN)
add_badge(slide, Inches(0.5), Inches(4.3), "SERVER DB", ACCENT_GREEN, Inches(1.2))
add_text(slide, Inches(1.8), Inches(4.3), Inches(4), Inches(0.3), "POS_Set[101] - S_Config", 14, TEXT_WHITE, True)
add_text(slide, Inches(0.5), Inches(4.7), Inches(5.8), Inches(0.3),
         "서버 동기화 업무규칙 / 쉼표구분 78+값 / S_Config_Call()로 로딩", 10, TEXT_GRAY)
items3 = [
    ("출력/인쇄 관련 (12)", "[0]Group_Print, [3]PrintLine, [15]Print_Barcode..."),
    ("판매/업무 규칙 (22)", "[4]Cash_Pass, [8]CancelAll, [14]No_Credit..."),
    ("포인트/회원 관련 (14)", "[6]Sale_Point, [25]Tran_Point, [18]CashBack..."),
    ("럭키유저/이벤트 (11)", "[34~44] LU_Chk, LU_Per, LU_Point..."),
    ("바코드/저울 (4)", "[31]Resend_Barcode, [56]Scale_18_YN..."),
    ("금액/세금/메시지 (15)", "[51]Change_Type, [62]Tax_Auto, [64]SaleMsg_YN..."),
]
for i, (cat, desc) in enumerate(items3):
    ry = Inches(5.05) + Inches(i * 0.3)
    add_text(slide, Inches(0.5), ry, Inches(1.8), Inches(0.28), cat, 9, ACCENT_GREEN, True)
    add_text(slide, Inches(2.3), ry, Inches(4.0), Inches(0.28), desc, 8, TEXT_LIGHT, False)

# Source 4: DB Office_User
add_shape(slide, Inches(6.8), Inches(4.2), Inches(6.2), Inches(2.9), border_color=ACCENT_ORANGE)
add_badge(slide, Inches(7.0), Inches(4.3), "SERVER DB", ACCENT_ORANGE, Inches(1.2))
add_text(slide, Inches(8.3), Inches(4.3), Inches(4), Inches(0.3), "Office_User - Shop", 14, TEXT_WHITE, True)
add_text(slide, Inches(7.0), Inches(4.7), Inches(5.8), Inches(0.3),
         "매장 기본 정보 / 서버DB 40+컬럼 / S_Config_Call()로 로딩", 10, TEXT_GRAY)
items4 = [
    ("매장 기본정보 (10)", "OFFICE_NAME, office_num, owner_name, address..."),
    ("SMS/푸시 설정 (6)", "eBill_sms_yn, eBill_push_yn, Push_title/msg..."),
    ("확장바코드 (1)", "strBarcode_YN (XBarcode)"),
    ("셀프 핫키 (5)", "selfPos_Hotkey1~4, SMS_GUBUN"),
    ("주방 메시지 (5)", "kitchenMsg1~5"),
]
for i, (cat, desc) in enumerate(items4):
    ry = Inches(5.05) + Inches(i * 0.3)
    add_text(slide, Inches(7.0), ry, Inches(1.8), Inches(0.28), cat, 9, ACCENT_ORANGE, True)
    add_text(slide, Inches(8.8), ry, Inches(4.0), Inches(0.28), desc, 8, TEXT_LIGHT, False)

add_text(slide, Inches(7.0), Inches(6.55), Inches(5.5), Inches(0.3),
         "로딩 순서: INI_Load() -> S_Config_Call() (DB 연결 후)", 9, TEXT_GRAY, False)


# ================================================================
# SLIDE 3: pos_config.ini - 공통 설정 (Application/Sale/Receipt/DataBase)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_BLUE)
add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
         "pos_config.ini - 공통 설정 (30키)", 22, TEXT_WHITE, True)
add_badge(slide, Inches(8.5), Inches(0.35), "INI / LOCAL", ACCENT_BLUE, Inches(1.3))

# [Application] - 5 keys
add_shape(slide, Inches(0.3), Inches(1.1), Inches(6.2), Inches(1.7), border_color=ACCENT_BLUE)
add_section_header(slide, Inches(0.5), Inches(1.2), "[Application] 애플리케이션 정보", 5, ACCENT_BLUE)
items_app = [
    ("Pro_Name", "프로그램명", "POSON"),
    ("Pro_Buttom", "하단 텍스트", "CopyRight..."),
    ("Font_Add", "폰트 추가", "1"),
    ("SALE_POINTUSE_UPDATE_YN", "포인트 업데이트 플래그", "0"),
    ("Backup_Month", "백업 주기(월)", "6"),
]
for i, (k, d, v) in enumerate(items_app):
    add_item_row(slide, Inches(0.5), Inches(1.6) + Inches(i * 0.22), k, d, v)

# [Sale] - 10 keys
add_shape(slide, Inches(0.3), Inches(3.0), Inches(6.2), Inches(2.8), border_color=ACCENT_GREEN)
add_section_header(slide, Inches(0.5), Inches(3.1), "[Sale] 판매 운영", 10, ACCENT_GREEN)
items_sale = [
    ("OpenDay", "영업 시작일", "날짜"),
    ("S_JeonPyo", "전표 시퀀스", "1"),
    ("St_Jeonpyo", "전표번호", "타임스탬프"),
    ("Befor_Tran", "이전거래 합계", "0"),
    ("Befor_Tran_Dec", "이전거래 차감", "0"),
    ("Befor_Tran_Card", "이전거래 카드", "0"),
    ("DataDown", "마스터 데이터 다운로드", "1"),
    ("Finish_Day", "마감일", "날짜"),
    ("Start_Price", "시재금", "0"),
    ("INOUT_Chk", "입출 체크", "1"),
]
for i, (k, d, v) in enumerate(items_sale):
    add_item_row(slide, Inches(0.5), Inches(3.5) + Inches(i * 0.22), k, d, v)

# [Receipt] - 9 keys
add_shape(slide, Inches(6.8), Inches(1.1), Inches(6.2), Inches(2.5), border_color=ACCENT_ORANGE)
add_section_header(slide, Inches(7.0), Inches(1.2), "[Receipt] 영수증 정보", 9, ACCENT_ORANGE)
items_rcpt = [
    ("ShopName", "매장명", "-"),
    ("ShopNumber", "사업자번호", "-"),
    ("Address", "주소", "-"),
    ("Owner", "대표자명", "-"),
    ("Tel", "전화번호", "-"),
    ("ShopName_Color", "매장명 색상", "8421440"),
    ("Top5", "상단 여백", "(빈값)"),
    ("Buttom5", "하단 여백", "(빈값)"),
    ("NewShopName_Color", "새매장명 색상", "16777215"),
]
for i, (k, d, v) in enumerate(items_rcpt):
    add_item_row(slide, Inches(7.0), Inches(1.6) + Inches(i * 0.22), k, d, v)

# [DataBase] - 6 keys
add_shape(slide, Inches(6.8), Inches(3.85), Inches(6.2), Inches(1.85), border_color=ACCENT_RED)
add_section_header(slide, Inches(7.0), Inches(3.95), "[DataBase] DB 연결", 6, ACCENT_RED)
items_db = [
    ("DBType", "DB 유형", "0"),
    ("IP", "DB 서버 IP", "192.168.10.48"),
    ("UserName", "DB 사용자 (암호화)", "-"),
    ("Pass", "DB 비밀번호 (암호화)", "-"),
    ("Port", "DB 포트", "1433"),
    ("DB_NAME", "DB 이름", "TIPS"),
]
for i, (k, d, v) in enumerate(items_db):
    add_item_row(slide, Inches(7.0), Inches(4.35) + Inches(i * 0.22), k, d, v)

add_text(slide, Inches(0.5), Inches(6.6), Inches(12), Inches(0.4),
         "참고: Sale 섹션은 날짜/시퀀스 등 기기별 운영 데이터  |  Receipt/Application은 매장 공통 설정", 10, TEXT_GRAY)


# ================================================================
# SLIDE 4: pos_config.ini - [Terminal] 터미널 HW (42키)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_TEAL)
add_text(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
         "pos_config.ini - [Terminal] 터미널 HW (42키)", 22, TEXT_WHITE, True)
add_badge(slide, Inches(8.2), Inches(0.35), "모두 기기별 설정", ACCENT_TEAL, Inches(1.5))

# Column 1: 기본/프린터 (14)
add_shape(slide, Inches(0.3), Inches(1.1), Inches(4.1), Inches(5.9), border_color=ACCENT_TEAL)
add_text(slide, Inches(0.5), Inches(1.2), Inches(3), Inches(0.3), "기본 / 프린터", 12, ACCENT_TEAL, True)
term_col1 = [
    ("TerminalType", "기기 종류"),
    ("PosNo", "POS 번호 (2자리)"),
    ("AdminPosNo", "관리자 POS번호"),
    ("Cashdraw", "캐시드로워"),
    ("Touch", "터치스크린"),
    ("Moniter", "메인 모니터"),
    ("S_Moniter", "서브 모니터"),
    ("Printer", "프린터 종류"),
    ("Printer_Port", "프린터 포트"),
    ("Printer_Serial_Bps", "프린터 전송속도"),
    ("Printer_CAT_USE", "CAT 프린터 사용"),
    ("CAT_Port", "CAT 포트"),
    ("CAT_BPS", "CAT 전송속도"),
    ("Printer_Bps", "프린터 속도 (중복)"),
]
for i, (k, d) in enumerate(term_col1):
    add_item_row(slide, Inches(0.5), Inches(1.55) + Inches(i * 0.2), k, d, size=8,
                 key_w=Inches(1.8), desc_w=Inches(2.0))

# Column 2: 스캐너/카드리더/CDP (14)
add_shape(slide, Inches(4.6), Inches(1.1), Inches(4.1), Inches(5.9), border_color=ACCENT_BLUE)
add_text(slide, Inches(4.8), Inches(1.2), Inches(3), Inches(0.3), "스캐너/카드리더/CDP", 12, ACCENT_BLUE, True)
term_col2 = [
    ("Scan_Name", "스캐너 종류"),
    ("Scan_Port", "스캐너 포트"),
    ("HandScan_Port", "핸드스캐너 포트"),
    ("MSR_Port", "카드리더 포트"),
    ("MSR_BPS", "카드리더 전송속도"),
    ("CDPName", "고객표시기명"),
    ("CDP_Port", "고객표시기 포트"),
    ("CDPLine", "고객표시기 라인"),
    ("CDPType", "고객표시기 종류"),
    ("CDP_BPS", "고객표시기 전송속도"),
    ("CDP_CASH_YN", "CDP 현금표시"),
    ("Coin_Name", "동전교환기 종류"),
    ("Coin_Port", "동전교환기 포트"),
    ("CBO_ScalePort", "저울 포트"),
]
for i, (k, d) in enumerate(term_col2):
    add_item_row(slide, Inches(4.8), Inches(1.55) + Inches(i * 0.2), k, d, size=8,
                 key_w=Inches(1.8), desc_w=Inches(2.0))

# Column 3: 주방/벨/듀얼 (14)
add_shape(slide, Inches(8.9), Inches(1.1), Inches(4.1), Inches(5.9), border_color=ACCENT_PURPLE)
add_text(slide, Inches(9.1), Inches(1.2), Inches(3), Inches(0.3), "주방프린터/벨/듀얼", 12, ACCENT_PURPLE, True)
term_col3 = [
    ("KitchenPrint", "주방프린터 종류"),
    ("KitchenPrintPort", "주방프린터 포트"),
    ("KitchenPrintBps", "주방프린터 전송속도"),
    ("Kitchen_fontsize", "주방 폰트크기"),
    ("SuPyo_Port", "수표 포트"),
    ("Dual", "듀얼 모니터"),
    ("Dual_Type", "듀얼 유형"),
    ("Bell_YN", "벨 사용"),
    ("Bell_Type", "벨 종류"),
    ("Bell_ComPort", "벨 COM포트"),
    ("Bell_Name", "벨 이름"),
    ("Bell_fID", "벨 시설ID"),
    ("Bell_fID_YN", "벨 시설ID 사용"),
    ("Bell_LEN", "벨 번호 길이"),
]
for i, (k, d) in enumerate(term_col3):
    add_item_row(slide, Inches(9.1), Inches(1.55) + Inches(i * 0.2), k, d, size=8,
                 key_w=Inches(1.8), desc_w=Inches(2.0))

add_text(slide, Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
         "전체 42키 모두 POS 기기별 하드웨어 설정 - TOBE DeviceSetting(TERMINAL 카테고리)로 이전", 9, TEXT_GRAY)


# ================================================================
# SLIDE 5: pos_config.ini - [Other] POS 업무 설정 (75키)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_GREEN)
add_text(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
         "pos_config.ini - [Other] POS 업무 설정 (75키)", 22, TEXT_WHITE, True)
add_badge(slide, Inches(9.5), Inches(0.35), "공통 + 기기별 혼재", ACCENT_GREEN, Inches(1.5))

# Column 1: 판매/가격 (25)
add_shape(slide, Inches(0.3), Inches(1.0), Inches(4.1), Inches(6.2), border_color=ACCENT_GREEN)
add_text(slide, Inches(0.5), Inches(1.1), Inches(3), Inches(0.3), "판매/가격 설정 (25키)", 11, ACCENT_GREEN, True)
col1 = [
    ("Price_Edit", "가격 수정 허용"),
    ("MaxPrice", "최대 결제금액"),
    ("HyeonGeumPrice", "최대 현금금액"),
    ("MIN_Card_Price", "최소 카드결제 금액"),
    ("SaleView", "판매화면 타입"),
    ("InPriView", "매입가 표시"),
    ("Grouping", "상품 그룹핑"),
    ("Grid_Fix_YN", "그리드 고정"),
    ("Grid_SaleEx", "판매 그리드 확장"),
    ("price_11_yn", "3단 가격 표시"),
    ("ENG_YN", "영어 모드"),
    ("Product_Sound", "상품 스캔 소리"),
    ("Product_CashOpen", "상품등록시 캐시드로워"),
    ("SCANCOP", "스캔 복사"),
    ("SCAN_REAL_Chk", "바코드 실시간 체크"),
    ("InfoDesk_YN", "안내데스크 모드"),
    ("InfoDesk_ViewAll", "안내데스크 전체보기"),
    ("Total_Qty_YN", "총 수량 표시"),
    ("Logo_Min_YN", "로고 최소화"),
    ("Free_Opt", "무료 옵션"),
    ("Grade_Memo_YN", "등급 메모"),
    ("Order_Call_YN", "주문 호출"),
    ("f_Boryu_YN", "보류 사용"),
    ("Boryu_Tran_Opt", "보류 거래 옵션"),
    ("Delay", "지연 설정"),
]
for i, (k, d) in enumerate(col1):
    add_item_row(slide, Inches(0.5), Inches(1.4) + Inches(i * 0.19), k, d, size=8,
                 key_w=Inches(1.9), desc_w=Inches(1.9))

# Column 2: 카드/결제/마감 (25)
add_shape(slide, Inches(4.6), Inches(1.0), Inches(4.1), Inches(6.2), border_color=ACCENT_BLUE)
add_text(slide, Inches(4.8), Inches(1.1), Inches(3), Inches(0.3), "카드/결제/마감 (25키)", 11, ACCENT_BLUE, True)
col2 = [
    ("OFF_CARD_Chk", "오프라인 카드 사용"),
    ("OFF_CARD_KEY_USE", "오프라인 카드키"),
    ("HAND_CARD_YN", "수동 카드입력"),
    ("eCard_YN", "전자카드 사용"),
    ("CARD_Timer_YN", "카드 타이머"),
    ("Card_Wav_Opt", "카드 결제 소리"),
    ("Card_View", "카드 표시"),
    ("GIFT_BILL_ETC", "상품권 출력 기타"),
    ("Gift_Input_YN", "상품권 입력"),
    ("SelfPos_YN", "셀프 모드 ON/OFF (!)"),
    ("MisuPrint", "미수 출력"),
    ("Re_Point", "환불시 포인트 재계산"),
    ("Re_Tax", "환불시 세금 재계산"),
    ("Re_CashBack", "환불시 캐시백 재계산"),
    ("MEMBER_ADD_SCREEN_YN", "회원 추가 화면"),
    ("noCVM_Bill_Print_YN", "비CVM 영수증 출력"),
    ("HOTKEY_ENTER_USE", "핫키 엔터 사용"),
    ("Touch_HotKey_Opt", "터치 핫키 옵션"),
    ("Day_F_Msg_Opt", "일마감 메시지 옵션"),
    ("Job_Finish_Cashdraw_YN", "마감시 캐시드로워"),
    ("All_Finish", "전체 마감"),
    ("Sale_Finish_Opt", "판매 마감 옵션"),
    ("Process_Method", "화면잠금 사용"),
    ("MDB_Compact_YN", "DB 자동 컴팩트"),
    ("Err_Write", "에러 로깅"),
]
for i, (k, d) in enumerate(col2):
    add_item_row(slide, Inches(4.8), Inches(1.4) + Inches(i * 0.19), k, d, size=8,
                 key_w=Inches(2.0), desc_w=Inches(1.8))

# Column 3: 출력/다운로드/기타 (25)
add_shape(slide, Inches(8.9), Inches(1.0), Inches(4.1), Inches(6.2), border_color=ACCENT_ORANGE)
add_text(slide, Inches(9.1), Inches(1.1), Inches(3), Inches(0.3), "출력/다운로드/기타 (25키)", 11, ACCENT_ORANGE, True)
col3 = [
    ("Print_VAT", "부가세 출력"),
    ("Print_Barcode", "바코드 출력"),
    ("Bott_Print", "하단 합계 출력"),
    ("Slot_Add", "용지 슬롯"),
    ("Cut_Position", "절단 위치"),
    ("Printer_OFF_Chk", "프린터 오프라인 체크"),
    ("Point_Bill_Print_YN", "포인트 영수증 출력"),
    ("reTranBill_Print_1_YN", "재거래 영수증 출력"),
    ("no_Bill_fMessage_YN", "무영수증 메시지"),
    ("no_Bill_fSound_YN", "무영수증 소리"),
    ("no_Bill_CusPoint_YN", "무영수증 포인트"),
    ("wSclae_Danwi_Print_YN", "저울 단위 출력"),
    ("wSclae_Danwi_View_YN", "저울 단위 표시"),
    ("mDown_YN", "마스터 다운로드 사용"),
    ("mDown_Week", "다운로드 주기(주)"),
    ("New_Item_Update", "신규상품 업데이트"),
    ("DataDown_Goods_NO", "상품 데이터 다운로드 번호"),
    ("Shop_Chk", "매장 체크"),
    ("Hotkey_Chk", "핫키 체크"),
    ("All_Data", "전체 데이터"),
    ("All_Finish_BAK", "전체마감 백업"),
    ("All_Finish_BAK_Def", "전체마감 백업 기본값"),
    ("PosMenu", "POS 메뉴 사용"),
    ("Dual_Pos_TEST", "듀얼POS 테스트"),
    ("Pro_Comp", "프로그램 압축"),
]
for i, (k, d) in enumerate(col3):
    add_item_row(slide, Inches(9.1), Inches(1.4) + Inches(i * 0.19), k, d, size=8,
                 key_w=Inches(2.0), desc_w=Inches(1.8))


# ================================================================
# SLIDE 6: pos_config.ini - VAN 결제 설정 (65+키)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_RED)
add_text(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
         "pos_config.ini - VAN 결제 설정 (13개 VAN사, 65+키)", 22, TEXT_WHITE, True)
add_badge(slide, Inches(9.5), Inches(0.35), "IP:공통 / 단말번호:기기별", ACCENT_RED, Inches(2.2))

# [Card] 공통
add_shape(slide, Inches(0.3), Inches(1.1), Inches(5.5), Inches(1.5), border_color=ACCENT_RED)
add_text(slide, Inches(0.5), Inches(1.2), Inches(4), Inches(0.3), "[Card] VAN 공통 설정", 13, ACCENT_RED, True)
card_items = [
    ("VAN_Select", "VAN사 선택 (1~13)", "기기별"),
    ("SingPad_Port", "서명패드 포트", "기기별"),
    ("CashBack_YN", "캐시백 사용", "공통"),
    ("OCash_Screen_YN", "온라인현금 화면", "공통"),
    ("Log_Delete", "로그 삭제", "공통"),
]
for i, (k, d, t) in enumerate(card_items):
    ry = Inches(1.55) + Inches(i * 0.2)
    add_text(slide, Inches(0.5), ry, Inches(1.7), Inches(0.2), k, 8, ACCENT_CYAN)
    add_text(slide, Inches(2.2), ry, Inches(2.2), Inches(0.2), d, 8, TEXT_LIGHT)
    tc = ACCENT_TEAL if t == "기기별" else ACCENT_ORANGE
    add_text(slide, Inches(4.5), ry, Inches(1.0), Inches(0.2), t, 8, tc, True)

# VAN사별 테이블
add_shape(slide, Inches(0.3), Inches(2.85), Inches(12.7), Inches(4.2), border_color=ACCENT_BLUE)
add_text(slide, Inches(0.5), Inches(2.95), Inches(4), Inches(0.3), "VAN사별 13개 섹션", 13, ACCENT_BLUE, True)

van_headers = ["VAN사", "INI 섹션", "IP (공통)", "Port (공통)", "DANMALNO (기기별)", "Snumber (기기별)", "기타 키"]
van_widths = [Inches(0.9), Inches(1.1), Inches(2.0), Inches(1.0), Inches(2.0), Inches(1.5), Inches(3.8)]
hx = Inches(0.5)
for h, w in zip(van_headers, van_widths):
    add_text(slide, hx, Inches(3.3), w, Inches(0.22), h, 9, ACCENT_BLUE, True)
    hx += w

van_data = [
    ("KSNET", "[KSNET]", "210.181.28.137", "9562", "O", "O", "SingPad_BPS, USB_YN"),
    ("KIS", "[KIS]", "-", "-", "O(암호화)", "O", "Pos_Number"),
    ("KMPS", "[KMPS]", "-", "-", "O", "O", "Port_chk"),
    ("SMARTRO", "[SMARTRO]", "211.192.50.244", "5801", "O", "O", "-"),
    ("StarVAN", "[StarVAN]", "-", "-", "O", "O", "-"),
    ("NICE", "[NICE]", "211.33.136.2", "7709", "O", "O", "Dll_Chk"),
    ("JTNet", "[JTNet]", "211.48.96.28", "11025", "O", "-", "-"),
    ("KICC", "[KICC]", "203.233.72.21", "15700", "O", "-", "SingPad_BPS, SingCall_Type"),
    ("KCB", "[KCB]", "-", "-", "O", "O", "Use"),
    ("KOVAN", "[KOVAN]", "203.231.12.152", "10902", "O(암호화)", "O", "Dll_Chk"),
    ("KOCES", "[KOCES]", "211.192.167.38", "10015", "O", "-", "SingPad_Name"),
    ("KCP", "[KCP]", "203.238.36.152", "9976", "O", "-", "Dll_Chk"),
    ("KFTC", "[KFTC]", "kftcvan.or.kr", "-", "O", "SeqNo", "Point"),
]
for i, row in enumerate(van_data):
    ry = Inches(3.58) + Inches(i * 0.25)
    rx = Inches(0.5)
    bg_c = BG_ROW if i % 2 == 0 else None
    if bg_c:
        add_shape(slide, Inches(0.4), ry - Inches(0.02), Inches(12.5), Inches(0.25), fill_color=bg_c)
    for j, (val, w) in enumerate(zip(row, van_widths)):
        c = TEXT_LIGHT
        if j == 0:
            c = ACCENT_ORANGE
        elif j in (2, 3):
            c = ACCENT_GREEN if val != "-" else TEXT_GRAY
        elif j in (4, 5):
            c = ACCENT_TEAL if "O" in val else TEXT_GRAY
        add_text(slide, rx, ry, w, Inches(0.22), val, 8, c, j == 0)
        rx += w

# 범례
add_shape(slide, Inches(6.0), Inches(1.1), Inches(7.0), Inches(1.5), border_color=BORDER_COLOR)
add_text(slide, Inches(6.2), Inches(1.2), Inches(6), Inches(0.3), "성격 분류 범례", 11, TEXT_WHITE, True)
legend = [
    ("IP / Port", "VAN사 제공 공통값, 매장별 동일", ACCENT_GREEN),
    ("DANMALNO / Snumber", "가맹점 단말기 고유번호, 기기별", ACCENT_TEAL),
    ("SingPad_BPS / SingCall_Type", "서명패드 HW 설정", ACCENT_ORANGE),
]
for i, (k, d, c) in enumerate(legend):
    ry = Inches(1.5) + Inches(i * 0.25)
    add_text(slide, Inches(6.2), ry, Inches(2.3), Inches(0.22), k, 9, c, True)
    add_text(slide, Inches(8.5), ry, Inches(4.0), Inches(0.22), d, 9, TEXT_LIGHT)


# ================================================================
# SLIDE 7: pos_config.ini - [Self] 셀프키오스크 전용 (73키)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_PINK)
add_text(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
         "pos_config.ini - [Self] 셀프키오스크 전용 (73키)", 22, TEXT_WHITE, True)
add_badge(slide, Inches(9.0), Inches(0.35), "Self_YN=1 일때 활성", ACCENT_PINK, Inches(2.0))

# Col 1: 현금결제(10) + 봉투/저울(6) + STL보안태그(6) = 22
add_shape(slide, Inches(0.3), Inches(1.0), Inches(4.1), Inches(6.2), border_color=ACCENT_PINK)
add_text(slide, Inches(0.5), Inches(1.1), Inches(3), Inches(0.25), "현금 결제 (10키)", 10, ACCENT_RED, True)
cash_items = [
    ("self_Cash", "현금 결제 사용"),
    ("self_CashPort", "현금기 COM포트"),
    ("self_CashSleep", "현금기 대기시간"),
    ("self_CashPhonNum", "현금기 연락처"),
    ("self_CashGubun", "현금기 종류"),
    ("self_NoHyunYoung", "실결제금액 숨김"),
    ("self_OneHPUse", "1만원권 사용"),
    ("self_50HPUse", "5만원권 사용"),
    ("self_10000Use", "만원권 사용"),
    ("Self_Card", "카드 결제 비활성화"),
]
for i, (k, d) in enumerate(cash_items):
    add_item_row(slide, Inches(0.5), Inches(1.35) + Inches(i * 0.19), k, d, size=8,
                 key_w=Inches(1.8), desc_w=Inches(1.9))

add_text(slide, Inches(0.5), Inches(3.3), Inches(3), Inches(0.25), "봉투/저울 (6키)", 10, ACCENT_TEAL, True)
bag_items = [
    ("self_BagPort", "봉투 감지 포트"),
    ("self_StartBag", "시작시 봉투 표시"),
    ("Self_MBagSell", "복수 봉투 판매"),
    ("self_LastBag", "마지막 봉투"),
    ("Self_ScalePort", "저울 COM포트"),
    ("Self_ScaleLimitG", "저울 무게 한도"),
]
for i, (k, d) in enumerate(bag_items):
    add_item_row(slide, Inches(0.5), Inches(3.55) + Inches(i * 0.19), k, d, size=8,
                 key_w=Inches(1.8), desc_w=Inches(1.9))

add_text(slide, Inches(0.5), Inches(4.75), Inches(3), Inches(0.25), "STL 보안태그 (6키)", 10, ACCENT_ORANGE, True)
stl_items = [
    ("STLGoods", "STL 상품 등록"),
    ("STLNoGoods", "STL 미등록 상품"),
    ("STLGoodsNo", "STL 상품번호"),
    ("STLSoundAdmin", "STL 관리자 알림음"),
    ("STLPort", "STL COM포트"),
    ("Bell", "셀프 벨"),
]
for i, (k, d) in enumerate(stl_items):
    add_item_row(slide, Inches(0.5), Inches(5.0) + Inches(i * 0.19), k, d, size=8,
                 key_w=Inches(1.8), desc_w=Inches(1.9))

# Col 2: 고객 인터페이스(13) + 자동운영(8) = 21
add_shape(slide, Inches(4.6), Inches(1.0), Inches(4.1), Inches(6.2), border_color=ACCENT_BLUE)
add_text(slide, Inches(4.8), Inches(1.1), Inches(3), Inches(0.25), "고객 인터페이스 (13키)", 10, ACCENT_BLUE, True)
cus_items = [
    ("SoundGuide", "음성 안내"),
    ("Self_CusNum4", "회원번호 4자리"),
    ("self_NoCustomer", "비회원 판매 허용"),
    ("self_CusSelect", "고객 선택 방식"),
    ("CusAddUse", "고객 추가 사용"),
    ("Self_CusAddEtc", "고객 추가 기타"),
    ("CusTopMsg", "고객 상단 메시지"),
    ("CusBTMsg1", "고객 버튼 메시지1"),
    ("CusBTMsg2", "고객 버튼 메시지2"),
    ("self_TouchSoundYN", "터치 소리"),
    ("self_CusAlarmUse", "고객 알람 사용"),
    ("Self_CusAlarmTime", "고객 알람 시간"),
    ("no_Bill_CusPoint_YN", "무영수증 고객포인트"),
]
for i, (k, d) in enumerate(cus_items):
    add_item_row(slide, Inches(4.8), Inches(1.35) + Inches(i * 0.19), k, d, size=8,
                 key_w=Inches(1.9), desc_w=Inches(1.8))

add_text(slide, Inches(4.8), Inches(3.95), Inches(3), Inches(0.25), "자동 운영 (8키)", 10, ACCENT_GREEN, True)
auto_items = [
    ("Auto_Open_YN", "자동 개점"),
    ("Auto_Finish_YN", "자동 마감"),
    ("Auto_Day", "자동 운영 날짜"),
    ("Auto_AP", "AM/PM"),
    ("Auto_HH", "시"),
    ("Auto_MM", "분"),
    ("Auto_ID", "자동 로그인 ID"),
    ("Auto_Pass", "자동 로그인 PW"),
]
for i, (k, d) in enumerate(auto_items):
    add_item_row(slide, Inches(4.8), Inches(4.2) + Inches(i * 0.19), k, d, size=8,
                 key_w=Inches(1.9), desc_w=Inches(1.8))

# Col 3: 포인트/알림(11) + UI(9) + 인쇄/기타(10) = 30
add_shape(slide, Inches(8.9), Inches(1.0), Inches(4.1), Inches(6.2), border_color=ACCENT_PURPLE)
add_text(slide, Inches(9.1), Inches(1.1), Inches(3), Inches(0.25), "포인트/알림 (11키)", 10, ACCENT_PURPLE, True)
pt_items = [
    ("self_NoAutoPoint", "자동포인트 비활성"),
    ("self_PointZero", "포인트 초기화"),
    ("self_PointHidden", "포인트 숨김"),
    ("Self_PointSMSUse", "포인트 SMS"),
    ("Self_UserCall", "직원호출"),
    ("Self_SMSAdmin", "관리자 SMS"),
    ("Self_Kakao", "카카오 알림"),
    ("Self_Zero", "제로클릭 알림"),
    ("self_ICSiren", "IC 사이렌"),
    ("Self_CamUse", "카메라 사용"),
    ("self_SNSGubun", "SNS 구분"),
]
for i, (k, d) in enumerate(pt_items):
    add_item_row(slide, Inches(9.1), Inches(1.35) + Inches(i * 0.19), k, d, size=8,
                 key_w=Inches(1.9), desc_w=Inches(1.8))

add_text(slide, Inches(9.1), Inches(3.5), Inches(3), Inches(0.25), "UI/화면 (9키)", 10, ACCENT_TEAL, True)
ui_items = [
    ("Self_MainPage", "메인페이지 표시"),
    ("Self_BTInit", "초기화 버튼 표시"),
    ("Self_OneCancel", "개별취소 버튼"),
    ("Self_zHotKey", "Z핫키 사용"),
    ("self_CountYN", "계수버튼 표시"),
    ("self_StartHotKey", "시작 핫키"),
    ("self_PriceUse", "가격 조정 사용"),
    ("self_PriceType", "가격 유형"),
    ("Self_Reader", "ID 리더기 유형"),
]
for i, (k, d) in enumerate(ui_items):
    add_item_row(slide, Inches(9.1), Inches(3.75) + Inches(i * 0.19), k, d, size=8,
                 key_w=Inches(1.9), desc_w=Inches(1.8))

add_text(slide, Inches(9.1), Inches(5.55), Inches(3), Inches(0.25), "인쇄/기타 (10키)", 10, ACCENT_ORANGE, True)
etc_items = [
    ("Self_AutoPrint", "자동 출력"),
    ("self_StoPrint", "출력 중지"),
    ("self_PrintAddress", "주소 출력"),
    ("self_PrintPhon", "전화번호 출력"),
    ("Self_JPYN", "일본 모드(동전교환)"),
    ("self_BagJPPort", "JP 봉투 포트"),
    ("Self_NoAutoGoods", "자동상품 비활성"),
    ("self_AppCard", "앱카드 사용"),
    ("self_Apple", "애플페이"),
    ("Self_Gif", "GIF 파일 경로"),
]
for i, (k, d) in enumerate(etc_items):
    add_item_row(slide, Inches(9.1), Inches(5.8) + Inches(i * 0.19), k, d, size=8,
                 key_w=Inches(1.9), desc_w=Inches(1.8))


# ================================================================
# SLIDE 8: 기타 기기 + POSON2 + UI/색상
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_PURPLE)
add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
         "pos_config.ini - 기타 기기 + POSON2 + UI/색상 (52키)", 22, TEXT_WHITE, True)

# FaceCam (4)
add_shape(slide, Inches(0.3), Inches(1.0), Inches(4.1), Inches(1.5), border_color=ACCENT_RED)
add_text(slide, Inches(0.5), Inches(1.1), Inches(3), Inches(0.25), "[FaceCam] 안면인식 (4키)", 11, ACCENT_RED, True)
fc_items = [("Use", "안면인식 사용"), ("TimeOut_Use", "타임아웃 사용"),
            ("Dealy", "지연시간 (초)"), ("TimeOut_Sec", "타임아웃 (초)")]
for i, (k, d) in enumerate(fc_items):
    add_item_row(slide, Inches(0.5), Inches(1.4) + Inches(i * 0.22), k, d, size=8,
                 key_w=Inches(1.5), desc_w=Inches(2.2))

# JaPan (7)
add_shape(slide, Inches(0.3), Inches(2.7), Inches(4.1), Inches(2.0), border_color=ACCENT_ORANGE)
add_text(slide, Inches(0.5), Inches(2.8), Inches(3), Inches(0.25), "[JaPan] 자판기 (7키)", 11, ACCENT_ORANGE, True)
jp_items = [("Auto_Open_YN", "자동 개점"), ("jp_Port1", "자판기 포트1"),
            ("jp_Port2", "자판기 포트2"), ("jp_StoYN", "재고관리"),
            ("jp_oneSell", "단일상품 판매"), ("jp_FirstPage", "첫페이지"),
            ("jp_SellPriShow", "판매가 표시")]
for i, (k, d) in enumerate(jp_items):
    add_item_row(slide, Inches(0.5), Inches(3.1) + Inches(i * 0.22), k, d, size=8,
                 key_w=Inches(1.5), desc_w=Inches(2.2))

# Selfimg (6) + Self21 (3)
add_shape(slide, Inches(0.3), Inches(4.9), Inches(4.1), Inches(2.4), border_color=ACCENT_TEAL)
add_text(slide, Inches(0.5), Inches(5.0), Inches(3), Inches(0.25), "[Selfimg] 이미지셀프 (6키)", 11, ACCENT_TEAL, True)
si_items = [("ScanUse", "스캔 사용"), ("onlyShop", "매장전용 판매"),
            ("prepayment", "선결제 사용"), ("selfImg_TableNoUse", "테이블번호 미사용"),
            ("selfImg_TableNumUse", "테이블숫자 사용"), ("self_SNSGubun", "SNS 구분")]
for i, (k, d) in enumerate(si_items):
    add_item_row(slide, Inches(0.5), Inches(5.3) + Inches(i * 0.2), k, d, size=8,
                 key_w=Inches(1.8), desc_w=Inches(2.0))

add_text(slide, Inches(0.5), Inches(6.55), Inches(3), Inches(0.25), "[Self21] 21인치셀프 (3키)", 11, ACCENT_PINK, True)
s21_items = [("Self21_CountYN", "계수 사용"), ("slef21_DCVisible", "할인 표시"), ("self_CusAddEtc", "고객추가 기타")]
for i, (k, d) in enumerate(s21_items):
    add_item_row(slide, Inches(0.5), Inches(6.8) + Inches(i * 0.2), k, d, size=8,
                 key_w=Inches(1.8), desc_w=Inches(2.0))

# Dual_Msg (7)
add_shape(slide, Inches(4.6), Inches(1.0), Inches(4.1), Inches(2.2), border_color=ACCENT_BLUE)
add_text(slide, Inches(4.8), Inches(1.1), Inches(3), Inches(0.25), "[Dual_Msg] 듀얼 모니터 (7키)", 11, ACCENT_BLUE, True)
dual_items = [
    ("DTop0", "듀얼 상단 메시지0"),
    ("DTop0_Color", "상단0 색상"),
    ("DTop1", "듀얼 상단 메시지1"),
    ("DTop1_Color", "상단1 색상"),
    ("DButtom0_Color", "하단 색상0"),
    ("DButtom1~3_Color", "하단 색상1~3"),
    ("DButtom4_Color", "하단 색상4"),
]
for i, (k, d) in enumerate(dual_items):
    add_item_row(slide, Inches(4.8), Inches(1.4) + Inches(i * 0.22), k, d, size=8,
                 key_w=Inches(1.8), desc_w=Inches(2.0))

# POSON2 (9)
add_shape(slide, Inches(4.6), Inches(3.45), Inches(4.1), Inches(2.8), border_color=ACCENT_GREEN)
add_text(slide, Inches(4.8), Inches(3.55), Inches(3), Inches(0.25), "[POSON2_*] 통합 기능 (9키)", 11, ACCENT_GREEN, True)
p2_items = [
    ("[UPSS] USE", "UPS 사용"),
    ("[UPSS] SMS_USE", "SMS 알림"),
    ("[UPSS] MSG_USE", "메시지 알림"),
    ("[UPSS] Recv_Mobile", "수신 번호"),
    ("[Login] USE", "자동로그인 사용"),
    ("[Login] sMoney", "시재금"),
    ("[Login] ID", "로그인 ID"),
    ("[wSock] USE_YN", "웹소켓 사용"),
    ("[wSock] TCP_PORT", "TCP 포트"),
]
for i, (k, d) in enumerate(p2_items):
    add_item_row(slide, Inches(4.8), Inches(3.85) + Inches(i * 0.22), k, d, size=8,
                 key_w=Inches(1.6), desc_w=Inches(2.2))

# Color/Font (16)
add_shape(slide, Inches(4.6), Inches(6.45), Inches(4.1), Inches(0.5), border_color=BORDER_COLOR)
add_text(slide, Inches(4.8), Inches(6.5), Inches(3.5), Inches(0.35),
         "[Sale_Color/Font] 색상/폰트 설정 (16키)", 10, TEXT_GRAY, True)

# 기타 INI 섹션 (18)
add_shape(slide, Inches(8.9), Inches(1.0), Inches(4.1), Inches(6.2), border_color=BORDER_COLOR)
add_text(slide, Inches(9.1), Inches(1.1), Inches(3.5), Inches(0.25), "기타 INI 섹션 (18키)", 11, TEXT_WHITE, True)
misc_items = [
    ("[INI_Pass] DelKey", "삭제 키"),
    ("[Backup] Path", "백업 경로"),
    ("[Update_Data] Convert_chk", "변환 체크"),
    ("[Option] Manual_YN", "수동 모드"),
    ("[MEMBER_ADD] Tax_Option", "세금 옵션"),
    ("[Product_Add] Update_Option", "업데이트 옵션"),
    ("[Product_Add] Customer_Code", "거래처 코드"),
    ("[Product_Add] Customer_Name", "거래처명"),
    ("[Product_Add] Group_Code", "그룹 코드"),
    ("[Product_Add] Group_Name", "그룹명"),
    ("[Product_Add] Branch_Rate", "분류 세율"),
    ("[Tag_Print] Size_Print", "태그 인쇄 크기"),
    ("[Tag_Print] Danwi_Print", "단위 인쇄"),
    ("[Tag_Print] NRP_Print", "NRP 인쇄"),
    ("[Tag_Print] Feed_Cnt", "피드 수"),
    ("[shop_chk] TT_yn", "매장 체크"),
    ("[Config] Screen_Hide_YN", "화면 숨김"),
    ("[Shop] KitchenMsg1~4", "주방 메시지"),
]
for i, (k, d) in enumerate(misc_items):
    add_item_row(slide, Inches(9.1), Inches(1.4) + Inches(i * 0.26), k, d, size=8,
                 key_w=Inches(2.2), desc_w=Inches(1.6))

add_text(slide, Inches(9.1), Inches(6.3), Inches(3.5), Inches(0.5),
         "[Grid_*] 20+섹션 (~100키)\nVB6 그리드 컬럼폭 - Vue/CSS로 대체", 9, TEXT_GRAY)


# ================================================================
# SLIDE 9: DB POS_Set[101] - S_Config (78+필드)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_GREEN)
add_text(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
         "DB: POS_Set[101] - S_Config (78+필드)", 22, TEXT_WHITE, True)
add_badge(slide, Inches(8.5), Inches(0.35), "SERVER DB / 매장 공통", ACCENT_GREEN, Inches(2.0))
add_text(slide, Inches(0.5), Inches(0.7), Inches(12), Inches(0.3),
         "서버 동기화 업무규칙 / 쉼표구분 78+값 / S_Config_Call()로 로딩", 10, TEXT_GRAY)

# Col 1: 출력/인쇄(12) + 판매/업무(11) = 23
add_shape(slide, Inches(0.3), Inches(1.1), Inches(4.1), Inches(6.1), border_color=ACCENT_GREEN)
add_text(slide, Inches(0.5), Inches(1.2), Inches(3), Inches(0.25), "출력/인쇄 관련 (12필드)", 10, ACCENT_GREEN, True)
pr_items = [
    ("[0] Group_Print", "그룹별 출력"),
    ("[3] PrintLine", "출력 라인(1/2줄)"),
    ("[11] Boryu_Print", "보류 출력"),
    ("[15] Print_Barcode", "바코드 출력"),
    ("[16] Print_VAT", "부가세 출력"),
    ("[17] Bott_Print", "하단 합계 출력"),
    ("[23] Card_Detail", "카드 상세 출력"),
    ("[24] SuPyo_Print", "수표 출력"),
    ("[28] Print_ShopBarcode", "매장바코드 출력"),
    ("[29] Print_SignJeonpyo", "서명패드 전표 출력"),
    ("[30] TranJeonpyo2_YN", "반품 전표 2매"),
    ("[53] Seq_Print", "순차 출력"),
]
for i, (k, d) in enumerate(pr_items):
    add_item_row(slide, Inches(0.5), Inches(1.5) + Inches(i * 0.19), k, d, size=8,
                 key_w=Inches(2.0), desc_w=Inches(1.8))

add_text(slide, Inches(0.5), Inches(3.85), Inches(3), Inches(0.25), "판매/업무 규칙 (11필드)", 10, ACCENT_BLUE, True)
sl_items = [
    ("[1] Money", "금액 표시"),
    ("[2] NewProduct", "신규상품"),
    ("[4] Cash_Pass", "현금 비밀번호"),
    ("[5] Price_Zero", "0원 판매 허용"),
    ("[8] CancelAll", "전체취소"),
    ("[9] JeonPyo", "전표 사용"),
    ("[10] Tran", "거래 사용"),
    ("[12] Card_CashOpen", "카드시 캐시드로워"),
    ("[13] CancelAll_Bigo", "전체취소 비고"),
    ("[14] No_Credit", "외상 불가"),
    ("[20] Grouping", "상품 그룹핑"),
]
for i, (k, d) in enumerate(sl_items):
    add_item_row(slide, Inches(0.5), Inches(4.15) + Inches(i * 0.19), k, d, size=8,
                 key_w=Inches(2.0), desc_w=Inches(1.8))

# Col 2: 판매/업무2(11) + 바코드(3) + 금액/세금(14) = 28
add_shape(slide, Inches(4.6), Inches(1.1), Inches(4.1), Inches(6.1), border_color=ACCENT_BLUE)
add_text(slide, Inches(4.8), Inches(1.2), Inches(3), Inches(0.25), "판매/업무 규칙2 (11필드)", 10, ACCENT_TEAL, True)
sl2_items = [
    ("[19] InPriView", "매입가 표시"),
    ("[21] BanPum_PassChk", "반품 비밀번호"),
    ("[22] MEMPrice_Chk", "회원 한도 체크"),
    ("[27] HyeonGeum_Chk", "현금 강제입력"),
    ("[32] Jeonpyo_HiddenYN", "전표번호 숨김"),
    ("[33] InTran_JeonpyoYN", "입금 전표 사용"),
    ("[46] Finish_TranChk", "마감 거래 불가"),
    ("[54] BanPum_Chk", "반품 전표 확인"),
    ("[55] BoryuMsg_View", "보류 메시지 표시"),
    ("[57] Boryu_Detail", "보류 상세 표시"),
    ("[65] Goods_Use", "판매시 재고 사용"),
]
for i, (k, d) in enumerate(sl2_items):
    add_item_row(slide, Inches(4.8), Inches(1.5) + Inches(i * 0.19), k, d, size=8,
                 key_w=Inches(2.0), desc_w=Inches(1.8))

add_text(slide, Inches(4.8), Inches(3.7), Inches(3), Inches(0.25), "바코드/저울 (3필드)", 10, ACCENT_ORANGE, True)
bc_items = [
    ("[31] Resend_BarcodeYN", "재전송 바코드"),
    ("[56] Scale_18_YN", "18자리 저울 바코드"),
    ("[102] Group_Sel", "그룹 선택(대/중)"),
]
for i, (k, d) in enumerate(bc_items):
    add_item_row(slide, Inches(4.8), Inches(4.0) + Inches(i * 0.19), k, d, size=8,
                 key_w=Inches(2.0), desc_w=Inches(1.8))

add_text(slide, Inches(4.8), Inches(4.65), Inches(3), Inches(0.25), "금액/세금/메시지 (15필드)", 10, ACCENT_RED, True)
amt_items = [
    ("[51] Change_Type", "거스름돈 유형"),
    ("[52] Change_Auto", "거스름돈 자동"),
    ("[58] Min_PrintYN", "최소금액 출력"),
    ("[59] Min_PointYN", "최소금액 포인트"),
    ("[60] Min_Price1", "최소금액1"),
    ("[61] Min_Price2", "최소금액2"),
    ("[62] Tax_Auto", "세금 자동발행"),
    ("[63] Profit_Msg", "이익 메시지"),
    ("[64] SaleMsg_YN", "판매 메시지"),
    ("[237] Tax_Gubun_View", "세금 구분 표시"),
    ("[241] POS_Smoney", "POS 준비금"),
    ("[242] FACE_ADD_Point_Chk", "안면인식 포인트 체크"),
    ("[243] FACE_ADD_Point", "안면인식 포인트"),
    ("[252~256] SalePrice_MSG*", "판매금액 메시지 (5종)"),
]
for i, (k, d) in enumerate(amt_items):
    add_item_row(slide, Inches(4.8), Inches(4.95) + Inches(i * 0.19), k, d, size=8,
                 key_w=Inches(2.1), desc_w=Inches(1.7))

# Col 3: 포인트/회원(14) + 럭키유저(11) = 25
add_shape(slide, Inches(8.9), Inches(1.1), Inches(4.1), Inches(6.1), border_color=ACCENT_PURPLE)
add_text(slide, Inches(9.1), Inches(1.2), Inches(3), Inches(0.25), "포인트/회원 관련 (14필드)", 10, ACCENT_PURPLE, True)
point_items = [
    ("[6] Sale_Point", "판매 포인트"),
    ("[7] Weight_Point", "중량 포인트"),
    ("[18] CashBack_YN", "캐시백 사용"),
    ("[25] Tran_Point", "거래 포인트"),
    ("[26] Sale_CashbackPoint", "캐시백 포인트 적립"),
    ("[45] Tran_Pointchk", "반품 포인트 체크"),
    ("[47] MEM_DecList_YN", "회원 차감 목록"),
    ("[48] MEM_DecList_Msg", "회원 차감 메시지"),
    ("[49] Dec_PointChk", "차감 포인트 체크"),
    ("[50] Dec_PointPer", "차감 포인트율"),
    ("[245] DEC_MEMO_Chk", "차감 메모 체크"),
    ("POS_Count", "POS 대수 (별도)"),
    ("Tran_PointPer", "거래 포인트율 (%)"),
    ("Dec_PointPer", "차감 포인트율 (%)"),
]
for i, (k, d) in enumerate(point_items):
    add_item_row(slide, Inches(9.1), Inches(1.5) + Inches(i * 0.19), k, d, size=8,
                 key_w=Inches(2.0), desc_w=Inches(1.8))

add_text(slide, Inches(9.1), Inches(4.25), Inches(3), Inches(0.25), "럭키유저/이벤트 (11필드)", 10, ACCENT_PINK, True)
lu_items = [
    ("[34] LU_Chk", "럭키유저 체크"),
    ("[35] LU_Per", "럭키유저 확률"),
    ("[36] LU_Point", "럭키유저 포인트"),
    ("[37] LU_CardPer", "럭키유저 카드 확률"),
    ("[38] LU_SPoint", "럭키유저 S포인트"),
    ("[39] LU_MinPrice", "럭키유저 최소금액"),
    ("[40] LU_Sound", "럭키유저 소리"),
    ("[41] LU_Print", "럭키유저 출력"),
    ("[42] LU_Opt1", "럭키유저 옵션1"),
    ("[43] LU_Opt2", "럭키유저 옵션2"),
    ("[44] LU_Price", "럭키유저 금액"),
]
for i, (k, d) in enumerate(lu_items):
    add_item_row(slide, Inches(9.1), Inches(4.55) + Inches(i * 0.19), k, d, size=8,
                 key_w=Inches(2.0), desc_w=Inches(1.8))


# ================================================================
# SLIDE 10: DB Office_User + 통계 요약 + 마이그레이션 맵
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_ORANGE)
add_text(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.6),
         "DB: Office_User (Shop) + 전체 통계 + 마이그레이션 맵", 22, TEXT_WHITE, True)

# Office_User 상세
add_shape(slide, Inches(0.3), Inches(1.0), Inches(5.8), Inches(3.2), border_color=ACCENT_ORANGE)
add_text(slide, Inches(0.5), Inches(1.1), Inches(5), Inches(0.3),
         "Office_User - Shop Type (24+필드)", 13, ACCENT_ORANGE, True)

ou_items = [
    ("OFFICE_NAME", "매장명", "공통"),
    ("office_num", "사업자번호", "공통"),
    ("owner_name", "대표자명", "공통"),
    ("office_tel1", "전화번호", "공통"),
    ("address1+2", "주소", "공통"),
    ("sto_cd", "매장코드", "공통"),
    ("version", "버전", "공통"),
    ("OFFICE_NAME2", "출력용 매장명", "공통"),
    ("Online_KEY", "온라인 키", "공통"),
    ("en_use", "고객 암호화", "공통"),
    ("eBill_sms_yn", "SMS 사용", "공통"),
    ("eBill_auto_sms_yn", "자동 SMS", "공통"),
    ("eBill_push_yn", "푸시 사용", "공통"),
    ("eBill_push_title", "푸시 제목", "공통"),
    ("eBill_push_msg", "푸시 메시지", "공통"),
    ("eBill_push_link", "푸시 링크", "공통"),
    ("strBarcode_YN", "확장바코드 사용", "공통"),
    ("selfPos_Hotkey1~4", "셀프 핫키 1~4", "셀프전용"),
    ("SMS_GUBUN", "SNS 구분", "셀프전용"),
    ("kitchenMsg1~5", "주방 메시지 1~5", "주방전용"),
]
for i, (k, d, scope) in enumerate(ou_items):
    ry = Inches(1.45) + Inches(i * 0.185)
    c = ACCENT_PINK if "셀프" in scope else (ACCENT_TEAL if "주방" in scope else ACCENT_CYAN)
    add_text(slide, Inches(0.5), ry, Inches(1.9), Inches(0.2), k, 8, c, False)
    add_text(slide, Inches(2.4), ry, Inches(2.0), Inches(0.2), d, 8, TEXT_LIGHT, False)
    if scope:
        add_text(slide, Inches(4.5), ry, Inches(1.2), Inches(0.2), scope, 7, TEXT_GRAY, False)

# 전체 통계
add_shape(slide, Inches(6.4), Inches(1.0), Inches(6.6), Inches(3.2), border_color=ACCENT_BLUE)
add_text(slide, Inches(6.6), Inches(1.1), Inches(6), Inches(0.3),
         "저장소별 항목 수 통계", 14, TEXT_WHITE, True)

stats = [
    ("pos_config.ini (POS 런타임)", "~446 키", ACCENT_BLUE),
    ("  [Other] POS 업무 설정", "75 키", TEXT_LIGHT),
    ("  [Terminal] HW 설정", "42 키", TEXT_LIGHT),
    ("  [Card]+VAN 13개사 결제", "65+ 키", TEXT_LIGHT),
    ("  [Self] 셀프키오스크 전용", "73 키", TEXT_LIGHT),
    ("  [FaceCam/JaPan/Selfimg/Self21/Dual/POSON2]", "52 키", TEXT_LIGHT),
    ("  [Grid_*] VB6 그리드 컬럼폭", "~100 키", TEXT_GRAY),
    ("  영수증/애플리케이션/DB접속/기타", "39 키", TEXT_LIGHT),
    ("", "", None),
    ("DB: POS_Set[101] (S_Config)", "~78 필드", ACCENT_GREEN),
    ("DB: Office_User (Shop)", "~24 필드", ACCENT_ORANGE),
    ("", "", None),
    ("POS 런타임 합계", "~548 항목", ACCENT_TEAL),
    ("config.ini (관리프로그램 전용)", "~350 키", ACCENT_PURPLE),
]
for i, (label, count, color) in enumerate(stats):
    if not label:
        continue
    ry = Inches(1.45) + Inches(i * 0.19)
    indent = Inches(6.8) if label.startswith("  ") else Inches(6.6)
    lbl = label.lstrip()
    add_text(slide, indent, ry, Inches(4.5), Inches(0.2), lbl, 9, color or TEXT_LIGHT, not label.startswith("  "))
    if count:
        add_text(slide, Inches(11.5), ry, Inches(1.3), Inches(0.2), count, 9, color or TEXT_LIGHT, True, PP_ALIGN.RIGHT)

# 마이그레이션 맵
add_shape(slide, Inches(0.3), Inches(4.45), Inches(12.7), Inches(2.8), border_color=ACCENT_TEAL)
add_text(slide, Inches(0.5), Inches(4.55), Inches(6), Inches(0.3),
         "ASIS -> TOBE 마이그레이션 맵", 14, ACCENT_TEAL, True)

mig_data = [
    ("ASIS 저장소", "ASIS 설정 유형", "항목 수", "TOBE 저장 위치", "비고"),
    ("INI [Other]+[Sale]+[Receipt]", "공통 설정 + 영수증 + 업무", "~100", "DB: SystemSetting", "6개 카테고리: SALE/PAYMENT/PRINT/POINT/BARCODE/SYSTEM"),
    ("INI [Terminal]", "터미널 HW 설정", "~42", "DB: DeviceSetting (TERMINAL)", "개별 POS/KIOSK/KITCHEN 기기별 관리"),
    ("INI [Card]+VAN", "VAN 결제 설정", "~65", "DB: DeviceSetting (VAN)", "12개 VAN사 통합 관리"),
    ("INI [Self]", "셀프키오스크 전용", "~73", "DB: DeviceSetting (SELF_*)", "6개 하위: CASH/BAG/AUTO/POINT/PRINT/ETC"),
    ("INI [FaceCam/JaPan/Selfimg/Self21]", "기타 기기 전용", "~20", "DB: DeviceSetting (TYPE_*)", "기기 유형별 분리 관리"),
    ("DB: POS_Set[101]", "업무 규칙 공통", "~78", "DB: SystemSetting", "기존 DB에서 DB로 이전 (구조 개선)"),
    ("DB: Office_User", "매장 기본 정보", "~24", "DB: SystemSetting / Business", "BusinessView.vue에서 관리"),
    ("INI [Grid_*]", "VB6 그리드 컬럼폭", "~100", "마이그레이션 제외 (Vue/CSS)", "대체 불필요 - Vue/CSS로 처리"),
    ("config.ini", "관리프로그램 전용", "~350", "별도 관리 프로그램", "POS 런타임과 무관"),
]

# Headers
hx = Inches(0.5)
h_widths = [Inches(3.0), Inches(1.8), Inches(0.8), Inches(2.5), Inches(4.5)]
for hi, (h_text, hw) in enumerate(zip(mig_data[0], h_widths)):
    add_text(slide, hx, Inches(4.9), hw, Inches(0.22), h_text, 9, ACCENT_TEAL, True)
    hx += hw

# Data rows
for i, row in enumerate(mig_data[1:]):
    ry = Inches(5.15) + Inches(i * 0.23)
    rx = Inches(0.5)
    bg_c = BG_ROW if i % 2 == 0 else None
    if bg_c:
        add_shape(slide, Inches(0.4), ry - Inches(0.01), Inches(12.5), Inches(0.23), fill_color=bg_c)
    colors = [ACCENT_CYAN, TEXT_LIGHT, ACCENT_ORANGE, ACCENT_GREEN, TEXT_GRAY]
    for j, (val, w) in enumerate(zip(row, h_widths)):
        bold = j == 0
        sz = 8 if j != 4 else 7
        add_text(slide, rx, ry, w, Inches(0.22), val, sz, colors[j], bold)
        rx += w


# Save
out_dir = os.path.join(os.path.dirname(__file__), "asis")
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "ASIS_INI_DB_\uc124\uc815_\ubd84\ub958_\uc804\uccb4\ubaa9\ub85d.pptx")
prs.save(out_path)
print("Saved: " + out_path)
print("Slides: " + str(len(prs.slides)))
