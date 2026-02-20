"""TOBE 매입~판매 비즈니스 플로우 PPTX 생성"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Color Palette ───
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1E, 0x29, 0x3B)
BG_ROW = RGBColor(0x16, 0x20, 0x33)
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)
ACCENT_GREEN = RGBColor(0x4A, 0xDE, 0x80)
ACCENT_ORANGE = RGBColor(0xFB, 0xBF, 0x24)
ACCENT_RED = RGBColor(0xF8, 0x71, 0x71)
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
ACCENT_TEAL = RGBColor(0x2D, 0xD4, 0xBF)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
TEXT_LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
BORDER_COLOR = RGBColor(0x33, 0x41, 0x55)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=14, color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.font.name = "Malgun Gothic"
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=12, color=TEXT_WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, text_color, is_bold, sz) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz or font_size)
        p.font.color.rgb = text_color or color
        p.font.bold = is_bold
        p.font.name = "Malgun Gothic"
        p.space_after = Pt(2)
    return txBox


def add_flow_arrow(slide, left, top, width=Inches(0.4), color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, left, top, height=Inches(0.35), color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, Inches(0.3), height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_badge(slide, left, top, text, color=ACCENT_BLUE, width=Inches(1.2)):
    shape = add_shape(slide, left, top, width, Inches(0.35), fill_color=color)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = TEXT_WHITE
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "Malgun Gothic"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


# ════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_TEAL)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
         "POSON POS  /  Electron + Vue 3 + Express", 20, TEXT_GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.3), Inches(11), Inches(1.2),
         "TOBE 매입 ~ 판매 비즈니스 플로우 분석", 42, TEXT_WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
         "Express + TypeScript + Prisma + PostgreSQL  |  Vue 3 + Pinia + Tailwind  |  Electron", 16, ACCENT_TEAL, False, PP_ALIGN.CENTER)

for i, (label, value) in enumerate([
    ("Tech Stack", "Express + Vue 3 + Electron"),
    ("DB", "PostgreSQL + Prisma ORM"),
    ("Cache", "Redis (5min TTL)"),
    ("분석일", "2026-02-19"),
]):
    x = Inches(1.5 + i * 2.7)
    card = add_shape(slide, x, Inches(5.0), Inches(2.4), Inches(1.1), fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text(slide, x + Inches(0.15), Inches(5.1), Inches(2.1), Inches(0.35), label, 10, TEXT_GRAY, False)
    add_text(slide, x + Inches(0.15), Inches(5.5), Inches(2.1), Inches(0.5), value, 13, TEXT_WHITE, True)


# ════════════════════════════════════════════════════════════════
# SLIDE 2: Business Cycle + DB Models
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_TEAL)
add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6), "전체 비즈니스 사이클 (TOBE)", 28, TEXT_WHITE, True)

cycle_items = [
    ("01", "거래처 등록", "Supplier\nSuppliersView", ACCENT_BLUE),
    ("02", "매입상품", "PurchaseProduct\n카탈로그+재고", ACCENT_GREEN),
    ("03", "매입 주문", "Purchase\n트랜잭션+StockMovement", ACCENT_ORANGE),
    ("04", "상품 등록", "Product\n다국어+M:N Category", ACCENT_TEAL),
    ("05", "판매/주문", "Order\nKiosk MenuView", ACCENT_RED),
    ("06", "결제", "PaymentService\nStrategy+Failover", ACCENT_PURPLE),
    ("07", "완료", "CompleteView\n30초 자동 복귀", ACCENT_BLUE),
]
for i, (num, title, desc, color) in enumerate(cycle_items):
    x = Inches(0.2 + i * 1.85)
    y = Inches(1.3)
    card = add_shape(slide, x, y, Inches(1.7), Inches(2.0), fill_color=BG_CARD, border_color=color)
    add_badge(slide, x + Inches(0.05), y + Inches(0.08), num, color, Inches(0.45))
    add_text(slide, x + Inches(0.1), y + Inches(0.5), Inches(1.5), Inches(0.4), title, 14, TEXT_WHITE, True)
    for j, line in enumerate(desc.split("\n")):
        add_text(slide, x + Inches(0.1), y + Inches(1.0 + j * 0.35), Inches(1.5), Inches(0.3), line, 9, TEXT_GRAY, False)
    if i < 6:
        add_flow_arrow(slide, x + Inches(1.7), y + Inches(0.85), Inches(0.15), color)

# DB Models
add_text(slide, Inches(0.6), Inches(3.6), Inches(12), Inches(0.4), "핵심 DB 모델 (PostgreSQL + Prisma)", 18, TEXT_WHITE, True)

models_left = [
    ("Supplier", "거래처", "id, code(S001~), name, type"),
    ("PurchaseProduct", "매입상품+재고", "id, barcode, stock, safeStock, costPrice"),
    ("Purchase", "매입 주문", "id, purchaseCode(P날짜-NNN), status"),
    ("PurchaseItem", "매입 상세", "purchaseId(CASCADE), qty, unitPrice"),
    ("StockMovement", "재고 이력", "type(6종), qty(±), before/after"),
]
models_right = [
    ("Category", "카테고리", "id, name(다국어4), imageUrl, sortOrder"),
    ("Product", "판매 상품", "id, barcode, sellPrice, status, options[]"),
    ("Order", "판매 주문", "id(UUID), orderNumber(YYMMDD-NNNN)"),
    ("OrderItem", "판매 상세", "orderId(CASCADE), productId, qty"),
    ("Payment", "결제", "id(UUID), vanCode, approvalNumber"),
]

for col, dataset in [(0, models_left), (1, models_right)]:
    x_base = Inches(0.4 + col * 6.5)
    for j, (model, desc, fields) in enumerate(dataset):
        y = Inches(4.15 + j * 0.55)
        if j % 2 == 0:
            add_shape(slide, x_base, y, Inches(6.2), Inches(0.5), fill_color=BG_ROW)
        add_text(slide, x_base + Inches(0.1), y + Inches(0.05), Inches(1.8), Inches(0.35), model, 11, ACCENT_TEAL, True)
        add_text(slide, x_base + Inches(2.0), y + Inches(0.05), Inches(1.2), Inches(0.35), desc, 10, TEXT_LIGHT, False)
        add_text(slide, x_base + Inches(3.3), y + Inches(0.05), Inches(2.8), Inches(0.35), fields, 9, TEXT_GRAY, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 3: Admin - 매입 Flow (UC-T01~T03)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_BLUE)
add_badge(slide, Inches(0.5), Inches(0.25), "ADMIN", ACCENT_BLUE, Inches(0.9))
add_text(slide, Inches(1.5), Inches(0.2), Inches(10), Inches(0.5), "매입 플로우: 거래처 → 매입상품 → 매입주문", 24, TEXT_WHITE, True)

# UC-T01
card1 = add_shape(slide, Inches(0.3), Inches(1.0), Inches(4.0), Inches(5.8), fill_color=BG_CARD, border_color=ACCENT_BLUE)
add_badge(slide, Inches(0.4), Inches(1.1), "UC-T01", ACCENT_BLUE, Inches(0.8))
add_text(slide, Inches(1.3), Inches(1.1), Inches(2.8), Inches(0.35), "거래처 등록", 16, TEXT_WHITE, True)
t01_steps = [
    ("SuppliersView 진입", TEXT_LIGHT),
    ("거래처 추가 버튼 클릭", TEXT_LIGHT),
    ("정보 입력 (이름 필수)", TEXT_LIGHT),
    ("POST /api/v1/suppliers", ACCENT_BLUE),
    ("사업자번호 중복 확인", TEXT_GRAY),
    ("코드 자동생성 (S001~)", ACCENT_GREEN),
    ("prisma.supplier.create()", ACCENT_GREEN),
    ("Redis 캐시 무효화", ACCENT_ORANGE),
]
for j, (step, color) in enumerate(t01_steps):
    add_text(slide, Inches(0.6), Inches(1.65 + j * 0.45), Inches(3.5), Inches(0.4), f"{j+1}. {step}", 11, color, False)
add_text(slide, Inches(0.6), Inches(5.5), Inches(3.5), Inches(0.4), "Supplier 모델", 10, TEXT_GRAY, True)
add_text(slide, Inches(0.6), Inches(5.85), Inches(3.5), Inches(0.8),
         "code(UNIQUE), name, type(5종)\nbusinessNumber, discountRate\nSoft Delete (isActive)", 9, TEXT_GRAY, False)

add_flow_arrow(slide, Inches(4.3), Inches(3.5), Inches(0.25), ACCENT_BLUE)

# UC-T02
card2 = add_shape(slide, Inches(4.6), Inches(1.0), Inches(4.0), Inches(5.8), fill_color=BG_CARD, border_color=ACCENT_GREEN)
add_badge(slide, Inches(4.7), Inches(1.1), "UC-T02", ACCENT_GREEN, Inches(0.8))
add_text(slide, Inches(5.6), Inches(1.1), Inches(2.8), Inches(0.35), "매입상품 등록", 16, TEXT_WHITE, True)
t02_steps = [
    ("PurchaseProductsView 진입", TEXT_LIGHT),
    ("바코드 스캔 or 자동생성", ACCENT_BLUE),
    ("  next-barcode API 호출", TEXT_GRAY),
    ("상품정보 입력 (6개 필수)", TEXT_LIGHT),
    ("거래처 선택 (supplierId)", TEXT_LIGHT),
    ("POST /purchase-products", ACCENT_BLUE),
    ("바코드 중복 확인", TEXT_GRAY),
    ("create (stock=0 초기)", ACCENT_GREEN),
]
for j, (step, color) in enumerate(t02_steps):
    add_text(slide, Inches(4.9), Inches(1.65 + j * 0.45), Inches(3.5), Inches(0.4), f"{'  ' if step.startswith('  ') else str(j+1) + '. '}{step.strip()}", 11, color, False)
add_text(slide, Inches(4.9), Inches(5.5), Inches(3.5), Inches(0.4), "PurchaseProduct 모델", 10, TEXT_GRAY, True)
add_text(slide, Inches(4.9), Inches(5.85), Inches(3.5), Inches(0.8),
         "barcode(UNIQUE), stock, safeStock\ncostPrice, vatAmount, taxType\nlCode/mCode/sCode (분류)", 9, TEXT_GRAY, False)

add_flow_arrow(slide, Inches(8.6), Inches(3.5), Inches(0.25), ACCENT_GREEN)

# UC-T03
card3 = add_shape(slide, Inches(8.9), Inches(1.0), Inches(4.1), Inches(5.8), fill_color=BG_CARD, border_color=ACCENT_ORANGE)
add_badge(slide, Inches(9.0), Inches(1.1), "UC-T03", ACCENT_ORANGE, Inches(0.8))
add_text(slide, Inches(9.9), Inches(1.1), Inches(2.9), Inches(0.35), "매입 주문", 16, TEXT_WHITE, True)
t03_steps = [
    ("PurchaseRegisterView 진입", TEXT_LIGHT),
    ("거래처 선택 (드롭다운)", TEXT_LIGHT),
    ("상품 검색 → 추가", TEXT_LIGHT),
    ("수량/단가 입력", TEXT_LIGHT),
    ("POST /api/v1/purchases", ACCENT_BLUE),
    ("$transaction 시작", ACCENT_ORANGE),
    ("  Purchase(CONFIRMED)", TEXT_GRAY),
    ("  stock += qty (증가)", ACCENT_GREEN),
    ("  StockMovement 기록", ACCENT_GREEN),
    ("캐시 무효화", ACCENT_ORANGE),
]
for j, (step, color) in enumerate(t03_steps):
    prefix = "  " if step.startswith("  ") else f"{j+1}. "
    add_text(slide, Inches(9.2), Inches(1.65 + j * 0.42), Inches(3.6), Inches(0.38), f"{prefix}{step.strip()}", 11, color, False)
add_text(slide, Inches(9.2), Inches(5.9), Inches(3.6), Inches(0.4), "코드: P{YYYYMMDD}-{NNN}", 10, ACCENT_ORANGE, True)
add_text(slide, Inches(9.2), Inches(6.25), Inches(3.6), Inches(0.4), "트랜잭션 필수 (재고 원자성)", 9, ACCENT_RED, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 4: Admin - 상품등록 + 재고 (UC-T04~T06)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_TEAL)
add_badge(slide, Inches(0.5), Inches(0.25), "ADMIN", ACCENT_BLUE, Inches(0.9))
add_text(slide, Inches(1.5), Inches(0.2), Inches(10), Inches(0.5), "상품 등록 → 재고 현황 → 매입 내역", 24, TEXT_WHITE, True)

# UC-T04 - Product Registration
card = add_shape(slide, Inches(0.3), Inches(1.0), Inches(4.0), Inches(6.0), fill_color=BG_CARD, border_color=ACCENT_TEAL)
add_badge(slide, Inches(0.4), Inches(1.1), "UC-T04", ACCENT_TEAL, Inches(0.8))
add_text(slide, Inches(1.3), Inches(1.1), Inches(2.8), Inches(0.35), "판매 상품 등록", 16, TEXT_WHITE, True)
t04_steps = [
    ("ProductsView 진입", TEXT_LIGHT),
    ("매입상품 연결 or 독립 등록", ACCENT_BLUE),
    ("  연결: barcode 자동 상속", TEXT_GRAY),
    ("상품정보 입력", TEXT_LIGHT),
    ("  name(다국어4), sellPrice", TEXT_GRAY),
    ("  categoryIds[] (M:N)", TEXT_GRAY),
    ("옵션 추가 (사이즈/토핑)", TEXT_LIGHT),
    ("POST /api/v1/products", ACCENT_BLUE),
    ("purchaseProductId 검증", TEXT_GRAY),
    ("categories M:N 연결", ACCENT_GREEN),
    ("캐시 무효화", ACCENT_ORANGE),
]
for j, (step, color) in enumerate(t04_steps):
    prefix = "  " if step.startswith("  ") else f"{j+1}. "
    add_text(slide, Inches(0.5), Inches(1.6 + j * 0.4), Inches(3.6), Inches(0.35), f"{prefix}{step.strip()}", 10, color, False)
add_text(slide, Inches(0.5), Inches(6.2), Inches(3.5), Inches(0.35), "상태: SELLING / SOLD_OUT / PENDING / HIDDEN", 9, ACCENT_TEAL, False)

add_flow_arrow(slide, Inches(4.3), Inches(3.5), Inches(0.25), ACCENT_TEAL)

# UC-T05 - Inventory Status
card = add_shape(slide, Inches(4.6), Inches(1.0), Inches(4.0), Inches(6.0), fill_color=BG_CARD, border_color=ACCENT_GREEN)
add_badge(slide, Inches(4.7), Inches(1.1), "UC-T05", ACCENT_GREEN, Inches(0.8))
add_text(slide, Inches(5.6), Inches(1.1), Inches(2.8), Inches(0.35), "재고 현황", 16, TEXT_WHITE, True)

summary_items = [
    ("전체 상품", "모든 매입상품 카운트", TEXT_LIGHT),
    ("정상 재고", "stock > safeStock", ACCENT_GREEN),
    ("낮은 재고", "0 < stock <= safeStock", ACCENT_ORANGE),
    ("품절", "stock <= 0", ACCENT_RED),
    ("재고 가치", "sum(stock x costPrice)", ACCENT_BLUE),
]
for j, (label, desc, color) in enumerate(summary_items):
    y = Inches(1.7 + j * 0.5)
    add_text(slide, Inches(4.9), y, Inches(1.5), Inches(0.35), label, 11, color, True)
    add_text(slide, Inches(6.4), y, Inches(2), Inches(0.35), desc, 10, TEXT_GRAY, False)

add_text(slide, Inches(4.9), Inches(4.3), Inches(3.5), Inches(0.35), "필터", 12, TEXT_WHITE, True)
filters = [
    "상품명/바코드 검색",
    "재고 상태: ALL / NORMAL / LOW / OUT",
    "재고 조정: POST /stock-movements/adjust",
    "재고 동기화: POST /sync-stock",
]
for j, f in enumerate(filters):
    add_text(slide, Inches(4.9), Inches(4.7 + j * 0.38), Inches(3.5), Inches(0.33), f"- {f}", 10, TEXT_LIGHT, False)

add_text(slide, Inches(4.9), Inches(6.2), Inches(3.5), Inches(0.35), "재고 변동: 매입(+), 판매(-), 조정(±)", 9, ACCENT_GREEN, False)

add_flow_arrow(slide, Inches(8.6), Inches(3.5), Inches(0.25), ACCENT_GREEN)

# UC-T06 - Purchase History
card = add_shape(slide, Inches(8.9), Inches(1.0), Inches(4.1), Inches(6.0), fill_color=BG_CARD, border_color=ACCENT_PURPLE)
add_badge(slide, Inches(9.0), Inches(1.1), "UC-T06", ACCENT_PURPLE, Inches(0.8))
add_text(slide, Inches(9.9), Inches(1.1), Inches(2.9), Inches(0.35), "매입 내역", 16, TEXT_WHITE, True)

add_text(slide, Inches(9.2), Inches(1.7), Inches(3.6), Inches(0.3), "3개 탭 구성", 12, TEXT_WHITE, True)
tabs = [
    ("날짜별", "날짜 범위 검색, 매입 목록", ACCENT_BLUE),
    ("거래처별", "거래처 그룹핑, 총액/건수", ACCENT_GREEN),
    ("결제", "미결제/기결제 구분", ACCENT_ORANGE),
]
for j, (tab, desc, color) in enumerate(tabs):
    y = Inches(2.15 + j * 0.65)
    add_badge(slide, Inches(9.3), y, tab, color, Inches(1.0))
    add_text(slide, Inches(10.4), y + Inches(0.02), Inches(2.5), Inches(0.3), desc, 10, TEXT_LIGHT, False)

add_text(slide, Inches(9.2), Inches(4.2), Inches(3.6), Inches(0.3), "필터", 12, TEXT_WHITE, True)
hist_filters = [
    "거래처 선택",
    "상태: DRAFT / CONFIRMED / CANCELLED",
    "날짜 범위",
    "상세 모달: 매입 품목 목록",
]
for j, f in enumerate(hist_filters):
    add_text(slide, Inches(9.2), Inches(4.55 + j * 0.35), Inches(3.6), Inches(0.3), f"- {f}", 10, TEXT_LIGHT, False)

add_text(slide, Inches(9.2), Inches(6.0), Inches(3.6), Inches(0.3), "매입 취소 시", 10, TEXT_WHITE, True)
add_text(slide, Inches(9.2), Inches(6.3), Inches(3.6), Inches(0.3), "DELETE → stock 자동 감소 (트랜잭션)", 9, ACCENT_RED, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 5: Kiosk Flow 1/2 (Welcome → Menu → Cart)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_GREEN)
add_badge(slide, Inches(0.5), Inches(0.25), "KIOSK", ACCENT_GREEN, Inches(0.9))
add_text(slide, Inches(1.5), Inches(0.2), Inches(10), Inches(0.5), "UC-K01: 키오스크 주문 플로우 (1/2) - 메뉴 선택", 24, TEXT_WHITE, True)

info_kiosk = [
    ("Actor", "고객 (셀프)"),
    ("View", "WelcomeView → MenuView"),
    ("Store", "useCartStore + useProductsStore"),
]
for i, (k, v) in enumerate(info_kiosk):
    x = Inches(0.5 + i * 4.2)
    add_text(slide, x, Inches(0.85), Inches(1), Inches(0.3), k, 10, TEXT_GRAY, True)
    add_text(slide, x + Inches(1.05), Inches(0.85), Inches(3), Inches(0.3), v, 11, ACCENT_GREEN, False)

# Phase cards
phases1 = [
    ("1단계", "환영 화면", "WelcomeView\n시작 버튼 터치\n언어 선택 (ko/en/ja/zh)", ACCENT_BLUE),
    ("2단계", "카테고리 선택", "좌측 사이드바\nproductsStore.selectCategory()\n프리셋 아이콘/이미지", ACCENT_GREEN),
    ("3단계", "상품 선택", "상품 그리드 (2-6컬럼)\n할인 배지 / 품절 오버레이\n재고 확인 (stock)", ACCENT_ORANGE),
    ("4단계", "옵션 선택", "옵션 모달 표시\n옵션 가격 합산\ncartStore.addItem()", ACCENT_PURPLE),
]
for i, (phase, title, desc, color) in enumerate(phases1):
    x = Inches(0.2 + i * 3.25)
    y = Inches(1.4)
    card = add_shape(slide, x, y, Inches(3.0), Inches(2.6), fill_color=BG_CARD, border_color=color)
    add_badge(slide, x + Inches(0.1), y + Inches(0.1), phase, color, Inches(0.9))
    add_text(slide, x + Inches(0.15), y + Inches(0.55), Inches(2.7), Inches(0.4), title, 16, TEXT_WHITE, True)
    for j, line in enumerate(desc.split("\n")):
        add_text(slide, x + Inches(0.15), y + Inches(1.1 + j * 0.4), Inches(2.7), Inches(0.35), line, 10, TEXT_GRAY, False)
    if i < 3:
        add_flow_arrow(slide, x + Inches(3.0), y + Inches(1.1), Inches(0.25), color)

# Cart section
add_text(slide, Inches(0.5), Inches(4.3), Inches(12), Inches(0.4), "장바구니 (MenuView 하단 푸터)", 18, TEXT_WHITE, True)
cart_card = add_shape(slide, Inches(0.3), Inches(4.8), Inches(12.7), Inches(2.3), fill_color=BG_CARD, border_color=ACCENT_BLUE)

cart_features = [
    ("주문 요약", "상품명 + 옵션 / 수량 ± 버튼 / 삭제 / 금액", ACCENT_BLUE),
    ("합계 표시", "총 개수 + 총 금액 (크고 굵게)", ACCENT_GREEN),
    ("결제하기", "빨간 버튼 (장바구니 비면 비활성화)", ACCENT_RED),
    ("취소", "cartStore.clear() → WelcomeView 복귀", TEXT_GRAY),
    ("토스트", "'상품명이 장바구니에 추가되었습니다' (2.5초)", ACCENT_ORANGE),
    ("품절 체크", "status=SOLD_OUT/PENDING or stock<=0 → 선택 불가", ACCENT_RED),
]
for j, (feat, desc, color) in enumerate(cart_features):
    col = j // 3
    row = j % 3
    x = Inches(0.6 + col * 6.3)
    y = Inches(5.0 + row * 0.55)
    add_text(slide, x, y, Inches(1.5), Inches(0.35), feat, 11, color, True)
    add_text(slide, x + Inches(1.6), y, Inches(4.5), Inches(0.35), desc, 10, TEXT_LIGHT, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 6: Kiosk Flow 2/2 (Payment → Complete)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_RED)
add_badge(slide, Inches(0.5), Inches(0.25), "KIOSK", ACCENT_GREEN, Inches(0.9))
add_text(slide, Inches(1.5), Inches(0.2), Inches(10), Inches(0.5), "UC-K01: 키오스크 주문 플로우 (2/2) - 결제 & 완료", 24, TEXT_WHITE, True)

# Payment flow
phases2 = [
    ("5단계", "결제 선택", "PaymentView\n카드/모바일/현금/스캐너\n주문 요약 + 세금(10%)", ACCENT_RED),
    ("6단계", "주문 생성", "POST /api/v1/orders\n주문번호: YYMMDD-NNNN\nstatus: PENDING", ACCENT_BLUE),
    ("7단계", "결제 처리", "POST /api/v1/payments\nVAN Strategy + Failover\nidempotencyKey 중복방지", ACCENT_ORANGE),
    ("8단계", "상태 변경", "PATCH /orders/:id/status\nPENDING → PAID\nstock -= qty (트랜잭션)", ACCENT_GREEN),
]
for i, (phase, title, desc, color) in enumerate(phases2):
    x = Inches(0.2 + i * 3.25)
    y = Inches(1.2)
    card = add_shape(slide, x, y, Inches(3.0), Inches(2.5), fill_color=BG_CARD, border_color=color)
    add_badge(slide, x + Inches(0.1), y + Inches(0.1), phase, color, Inches(0.9))
    add_text(slide, x + Inches(0.15), y + Inches(0.55), Inches(2.7), Inches(0.4), title, 16, TEXT_WHITE, True)
    for j, line in enumerate(desc.split("\n")):
        add_text(slide, x + Inches(0.15), y + Inches(1.1 + j * 0.4), Inches(2.7), Inches(0.35), line, 10, TEXT_GRAY, False)
    if i < 3:
        add_flow_arrow(slide, x + Inches(3.0), y + Inches(1.0), Inches(0.25), color)

# Complete & Return
phases3 = [
    ("9단계", "주문 완료", "CompleteView\n주문번호 표시\n승인번호 or 거스름돈", ACCENT_PURPLE),
    ("10단계", "카운트다운", "30초 자동 카운트\n홈으로 돌아가기 버튼\nlocaleStore.resetLocale()", ACCENT_BLUE),
    ("11단계", "홈 복귀", "WelcomeView로 이동\ncartStore.clear()\n언어 초기화 → 대기 상태", ACCENT_GREEN),
]
for i, (phase, title, desc, color) in enumerate(phases3):
    x = Inches(0.2 + i * 3.25)
    y = Inches(4.0)
    card = add_shape(slide, x, y, Inches(3.0), Inches(2.3), fill_color=BG_CARD, border_color=color)
    add_badge(slide, x + Inches(0.1), y + Inches(0.1), phase, color, Inches(0.9))
    add_text(slide, x + Inches(0.15), y + Inches(0.55), Inches(2.7), Inches(0.4), title, 16, TEXT_WHITE, True)
    for j, line in enumerate(desc.split("\n")):
        add_text(slide, x + Inches(0.15), y + Inches(1.1 + j * 0.38), Inches(2.7), Inches(0.33), line, 10, TEXT_GRAY, False)
    if i < 2:
        add_flow_arrow(slide, x + Inches(3.0), y + Inches(1.0), Inches(0.25), color)

# Alternative Flow box
alt_card = add_shape(slide, Inches(10.0), Inches(4.0), Inches(3.1), Inches(2.8), fill_color=BG_CARD, border_color=ACCENT_RED)
add_text(slide, Inches(10.15), Inches(4.1), Inches(2.8), Inches(0.35), "Alternative Flow", 12, TEXT_WHITE, True)
alt_items = [
    ("5a", "오프라인+카드 → 불가", ACCENT_RED),
    ("7a", "VAN 실패 → Failover", ACCENT_ORANGE),
    ("7b", "CircuitBreaker OPEN", ACCENT_ORANGE),
    ("7c", "중복결제 → 기존 반환", TEXT_LIGHT),
    ("8a", "취소 → 재고 복구", ACCENT_RED),
    ("10a", "수동 홈 복귀", TEXT_LIGHT),
]
for j, (code, desc, color) in enumerate(alt_items):
    y = Inches(4.55 + j * 0.35)
    add_text(slide, Inches(10.2), y, Inches(0.4), Inches(0.3), code, 9, color, True)
    add_text(slide, Inches(10.65), y, Inches(2.3), Inches(0.3), desc, 9, TEXT_GRAY, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 7: State Machines & StockMovement
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_PURPLE)
add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6), "상태 머신 & 재고 이동 추적", 28, TEXT_WHITE, True)

# Order Status Machine
add_text(slide, Inches(0.5), Inches(1.1), Inches(6), Inches(0.4), "주문 상태 (Order Status)", 16, TEXT_WHITE, True)
order_card = add_shape(slide, Inches(0.3), Inches(1.6), Inches(6.2), Inches(3.0), fill_color=BG_CARD, border_color=ACCENT_BLUE)
order_states = [
    ("PENDING", "주문 생성 (재고 미차감)", TEXT_LIGHT, ACCENT_BLUE),
    ("  ↓  결제 완료", "", TEXT_GRAY, None),
    ("PAID", "재고 차감: deductStockForOrder()", ACCENT_GREEN, ACCENT_GREEN),
    ("  ↓  준비 시작", "", TEXT_GRAY, None),
    ("PREPARING", "조리/준비 중", TEXT_LIGHT, ACCENT_ORANGE),
    ("  ↓  완료 or 취소", "", TEXT_GRAY, None),
    ("COMPLETED", "completedAt 설정", ACCENT_GREEN, ACCENT_GREEN),
    ("CANCELLED", "재고 복구: restoreStockForOrder()", ACCENT_RED, ACCENT_RED),
]
for j, (state, desc, color, badge_color) in enumerate(order_states):
    y = Inches(1.75 + j * 0.35)
    if badge_color:
        add_badge(slide, Inches(0.5), y, state, badge_color, Inches(1.5))
        add_text(slide, Inches(2.2), y + Inches(0.02), Inches(4), Inches(0.3), desc, 10, color, False)
    else:
        add_text(slide, Inches(0.8), y, Inches(5.5), Inches(0.3), state, 10, color, False)

# Purchase Status Machine
add_text(slide, Inches(0.5), Inches(4.9), Inches(6), Inches(0.4), "매입 상태 (Purchase Status)", 16, TEXT_WHITE, True)
purchase_card = add_shape(slide, Inches(0.3), Inches(5.35), Inches(6.2), Inches(1.7), fill_color=BG_CARD, border_color=ACCENT_ORANGE)
p_states = [
    ("DRAFT", "임시저장 (재고 미변동)", TEXT_LIGHT),
    ("CONFIRMED", "확정 → stock += qty, PURCHASE_IN", ACCENT_GREEN),
    ("CANCELLED", "취소 → stock -= qty, PURCHASE_CANCEL", ACCENT_RED),
]
for j, (state, desc, color) in enumerate(p_states):
    y = Inches(5.5 + j * 0.45)
    add_badge(slide, Inches(0.5), y, state, color if color != TEXT_LIGHT else ACCENT_BLUE, Inches(1.5))
    add_text(slide, Inches(2.2), y + Inches(0.02), Inches(4), Inches(0.3), desc, 10, color, False)

# StockMovement
add_text(slide, Inches(7.0), Inches(1.1), Inches(6), Inches(0.4), "재고 이동 추적 (StockMovement)", 16, TEXT_WHITE, True)
stock_card = add_shape(slide, Inches(6.8), Inches(1.6), Inches(6.2), Inches(5.45), fill_color=BG_CARD, border_color=ACCENT_TEAL)

sm_types = [
    ("PURCHASE_IN", "+", "매입 확정", "purchases.ts POST", ACCENT_GREEN),
    ("PURCHASE_CANCEL", "-", "매입 취소/수정", "purchases.ts DEL/PATCH", ACCENT_RED),
    ("SALE_OUT", "-", "주문 결제 (PAID)", "orders.ts deductStock", ACCENT_RED),
    ("SALE_CANCEL", "+", "주문 취소", "orders.ts restoreStock", ACCENT_GREEN),
    ("ADJUSTMENT", "±", "재고 수동 조정", "stock-movements.ts adjust", ACCENT_ORANGE),
    ("SYNC", "±", "재고 동기화", "purchase-products.ts sync", ACCENT_PURPLE),
]

# Header
add_text(slide, Inches(7.0), Inches(1.75), Inches(1.8), Inches(0.3), "Type", 10, TEXT_GRAY, True)
add_text(slide, Inches(8.8), Inches(1.75), Inches(0.4), Inches(0.3), "±", 10, TEXT_GRAY, True)
add_text(slide, Inches(9.2), Inches(1.75), Inches(1.5), Inches(0.3), "트리거", 10, TEXT_GRAY, True)
add_text(slide, Inches(10.8), Inches(1.75), Inches(2), Inches(0.3), "호출 위치", 10, TEXT_GRAY, True)

for j, (sm_type, sign, trigger, location, color) in enumerate(sm_types):
    y = Inches(2.2 + j * 0.55)
    if j % 2 == 0:
        add_shape(slide, Inches(6.9), y - Inches(0.05), Inches(6.0), Inches(0.5), fill_color=BG_ROW)
    add_text(slide, Inches(7.0), y, Inches(1.8), Inches(0.35), sm_type, 10, color, True)
    add_text(slide, Inches(8.8), y, Inches(0.4), Inches(0.35), sign, 12, color, True)
    add_text(slide, Inches(9.2), y, Inches(1.5), Inches(0.35), trigger, 10, TEXT_LIGHT, False)
    add_text(slide, Inches(10.8), y, Inches(2), Inches(0.35), location, 9, TEXT_GRAY, False)

# Record format
add_text(slide, Inches(7.0), Inches(5.6), Inches(5.8), Inches(0.3), "기록 필드: type, quantity(±), stockBefore, stockAfter", 10, ACCENT_TEAL, False)
add_text(slide, Inches(7.0), Inches(5.95), Inches(5.8), Inches(0.3), "추적: purchaseId / orderId / adjustmentCode / reason / createdBy", 9, TEXT_GRAY, False)
add_text(slide, Inches(7.0), Inches(6.3), Inches(5.8), Inches(0.3), "조정 코드: ADJ-YYYYMMDD-NNNN (일괄 조정 그룹)", 9, TEXT_GRAY, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 8: Payment System Architecture
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_RED)
add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6), "결제 시스템 아키텍처", 28, TEXT_WHITE, True)

# Strategy Pattern
add_text(slide, Inches(0.5), Inches(1.1), Inches(6), Inches(0.4), "Strategy Pattern + CircuitBreaker", 16, TEXT_WHITE, True)
strat_card = add_shape(slide, Inches(0.3), Inches(1.6), Inches(6.2), Inches(5.4), fill_color=BG_CARD, border_color=ACCENT_RED)

# PaymentService box
add_shape(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(0.5), fill_color=ACCENT_RED)
add_text(slide, Inches(0.7), Inches(1.85), Inches(5.4), Inches(0.4), "PaymentService", 14, TEXT_WHITE, True, PP_ALIGN.CENTER)

add_down_arrow(slide, Inches(3.2), Inches(2.3), Inches(0.3), ACCENT_RED)

# VAN Strategies
vans = [
    ("NICE", "NicePaymentStrategy", ACCENT_BLUE),
    ("KICC", "KiccPaymentStrategy", ACCENT_GREEN),
    ("KIS", "KisPaymentStrategy", ACCENT_ORANGE),
    ("SMARTRO", "SmartroPaymentStrategy", ACCENT_PURPLE),
]
for i, (van, strategy, color) in enumerate(vans):
    x = Inches(0.5 + i * 1.45)
    y = Inches(2.7)
    add_shape(slide, x, y, Inches(1.35), Inches(0.8), fill_color=color)
    add_text(slide, x + Inches(0.05), y + Inches(0.05), Inches(1.25), Inches(0.3), van, 11, TEXT_WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.05), y + Inches(0.35), Inches(1.25), Inches(0.3), "Strategy", 8, RGBColor(0xFF, 0xFF, 0xFF), False, PP_ALIGN.CENTER)

# CircuitBreaker
add_text(slide, Inches(0.7), Inches(3.8), Inches(5.5), Inches(0.3), "CircuitBreaker 상태", 12, TEXT_WHITE, True)
cb_states = [
    ("CLOSED", "정상 → 요청 통과", ACCENT_GREEN),
    ("OPEN", "장애 → 요청 차단, Failover", ACCENT_RED),
    ("HALF_OPEN", "복구 시도 → 제한된 요청", ACCENT_ORANGE),
]
for j, (state, desc, color) in enumerate(cb_states):
    y = Inches(4.2 + j * 0.45)
    add_badge(slide, Inches(0.7), y, state, color, Inches(1.3))
    add_text(slide, Inches(2.2), y + Inches(0.02), Inches(4), Inches(0.3), desc, 10, TEXT_LIGHT, False)

# Failover
add_text(slide, Inches(0.7), Inches(5.7), Inches(5.5), Inches(0.3), "Failover 순서", 12, TEXT_WHITE, True)
add_text(slide, Inches(0.7), Inches(6.05), Inches(5.5), Inches(0.3),
         "Primary: NICE → Backup: KICC → KIS", 11, ACCENT_BLUE, True)
add_text(slide, Inches(0.7), Inches(6.4), Inches(5.5), Inches(0.3),
         "SUCCESS → 반환 / RETRYABLE → 다음 VAN / FAILURE → 에러", 10, TEXT_GRAY, False)

# Idempotency & Features
add_text(slide, Inches(7.0), Inches(1.1), Inches(6), Inches(0.4), "안전 장치 & 특징", 16, TEXT_WHITE, True)
feat_card = add_shape(slide, Inches(6.8), Inches(1.6), Inches(6.2), Inches(5.4), fill_color=BG_CARD, border_color=ACCENT_PURPLE)

features = [
    ("멱등성 (Idempotency)", "idempotencyKey로 중복 결제 완벽 방지\n동일 키 → 기존 결과 반환 (재처리 없음)", ACCENT_BLUE),
    ("결제 수단", "신용카드: VAN Strategy 패턴\n모바일 페이: 카카오페이/네이버페이\n현금: 로컬 처리 (거스름돈 계산)\n스캐너: 준비 중", ACCENT_GREEN),
    ("오프라인 대응", "현금 결제만 가능\n카드/모바일 → '오프라인 상태입니다' 경고\n주문 로컬 저장 → 온라인 시 동기화", ACCENT_ORANGE),
    ("트랜잭션 보장", "주문 생성: 주문번호 race condition 방지\n상태 변경: 재고 처리 원자성\n결제 실패 시: 주문 상태 롤백", ACCENT_RED),
]
for j, (title, desc, color) in enumerate(features):
    y = Inches(1.8 + j * 1.25)
    add_badge(slide, Inches(7.0), y, title, color, Inches(2.5))
    for k, line in enumerate(desc.split("\n")):
        add_text(slide, Inches(7.1), y + Inches(0.4 + k * 0.28), Inches(5.7), Inches(0.25), line, 10, TEXT_LIGHT, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 9: ASIS vs TOBE Comparison
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_PURPLE)
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "ASIS vs TOBE 아키텍처 비교", 28, TEXT_WHITE, True)

comparison = [
    ("항목", "ASIS (VB6)", "TOBE (Express + Vue 3)"),
    ("DB", "MS-SQL, 월별 동적 테이블", "PostgreSQL + Prisma ORM, 단일 테이블"),
    ("재고", "LastSt_{YYYYMM} 이월재고", "PurchaseProduct.stock + StockMovement"),
    ("코드 생성", "사용자 수동 입력", "자동 (S001, P날짜-NNN, YYMMDD-NNNN)"),
    ("VAN", "12사 DLL 직접 호출", "4사 Strategy + CircuitBreaker + Failover"),
    ("결제 안전", "N/A", "idempotencyKey 멱등성 보장"),
    ("모드 전환", "C_Config.Self_YN", "Device 모델 (POS/KIOSK/KITCHEN)"),
    ("다국어", "N/A", "4개 언어 (ko/en/ja/zh) DB + i18n"),
    ("캐시", "N/A", "Redis 5분 TTL + CUD 무효화"),
    ("오프라인", "Access MDB", "IndexedDB + SyncManager"),
    ("트랜잭션", "ADO 수동 관리", "Prisma $transaction() 자동"),
    ("상태 추적", "단순 플래그", "상태 머신 (5단계 주문, 3단계 매입)"),
    ("재고 추적", "단순 증감", "StockMovement 6타입 (감사 추적)"),
    ("인증", "Admin_Gubun 레벨", "JWT + 역할 (4단계 RBAC)"),
    ("설정", "INI + DB POS_Set", "2계층: SystemSetting + DeviceSetting"),
]

# Header
header_y = Inches(1.1)
add_shape(slide, Inches(0.3), header_y, Inches(2.2), Inches(0.45), fill_color=RGBColor(0x33, 0x41, 0x55))
add_shape(slide, Inches(2.6), header_y, Inches(5.0), Inches(0.45), fill_color=ACCENT_RED)
add_shape(slide, Inches(7.7), header_y, Inches(5.3), Inches(0.45), fill_color=ACCENT_TEAL)

for j, (item, asis, tobe) in enumerate(comparison):
    y = header_y + Inches(j * 0.4)
    if j == 0:
        add_text(slide, Inches(0.5), y + Inches(0.05), Inches(2), Inches(0.35), item, 11, TEXT_WHITE, True)
        add_text(slide, Inches(2.8), y + Inches(0.05), Inches(4.7), Inches(0.35), asis, 11, TEXT_WHITE, True)
        add_text(slide, Inches(7.9), y + Inches(0.05), Inches(5), Inches(0.35), tobe, 11, TEXT_WHITE, True)
    else:
        if j % 2 == 0:
            add_shape(slide, Inches(0.3), y, Inches(12.7), Inches(0.4), fill_color=BG_ROW)
        add_text(slide, Inches(0.5), y + Inches(0.03), Inches(2), Inches(0.35), item, 10, TEXT_GRAY, True)
        add_text(slide, Inches(2.8), y + Inches(0.03), Inches(4.7), Inches(0.35), asis, 10, ACCENT_RED, False)
        add_text(slide, Inches(7.9), y + Inches(0.03), Inches(5), Inches(0.35), tobe, 10, ACCENT_TEAL, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 10: Integrated Scenario (8 Steps)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_TEAL)
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "통합 시나리오: Purchase-to-Sale 완전한 8단계", 28, TEXT_WHITE, True)

scenarios = [
    ("1", "거래처 등록", "POST /suppliers\n→ S001", ACCENT_BLUE, "-"),
    ("2", "매입상품", "POST /purchase-products\n→ stock=0", ACCENT_GREEN, "-"),
    ("3", "매입 주문", "POST /purchases\n[트랜잭션]", ACCENT_ORANGE, "stock\n+50"),
    ("4", "상품 등록", "POST /products\n→ M:N Category", ACCENT_TEAL, "-"),
    ("5", "키오스크 주문", "POST /orders\n→ PENDING", ACCENT_RED, "-"),
    ("6", "결제", "POST /payments\n→ APPROVED", ACCENT_PURPLE, "-"),
    ("7", "상태 PAID", "PATCH /orders/status\n[트랜잭션]", ACCENT_GREEN, "stock\n-3"),
    ("8", "완료", "COMPLETED\n30초 → 홈", ACCENT_BLUE, "-"),
]

for i, (num, title, desc, color, stock) in enumerate(scenarios):
    col = i % 4
    row = i // 4
    x = Inches(0.2 + col * 3.25)
    y = Inches(1.2 + row * 3.2)
    card = add_shape(slide, x, y, Inches(3.0), Inches(2.8), fill_color=BG_CARD, border_color=color)

    # Number circle
    add_badge(slide, x + Inches(0.1), y + Inches(0.1), num, color, Inches(0.4))
    add_text(slide, x + Inches(0.6), y + Inches(0.1), Inches(2.2), Inches(0.35), title, 15, TEXT_WHITE, True)

    for j, line in enumerate(desc.split("\n")):
        accent = ACCENT_ORANGE if "[트랜잭션]" in line else ACCENT_BLUE if "POST" in line or "PATCH" in line else TEXT_LIGHT
        add_text(slide, x + Inches(0.15), y + Inches(0.6 + j * 0.35), Inches(2.7), Inches(0.3), line, 11, accent, False)

    # Stock impact
    if stock != "-":
        stock_color = ACCENT_GREEN if "+" in stock else ACCENT_RED
        add_badge(slide, x + Inches(1.8), y + Inches(1.4), stock.replace("\n", " "), stock_color, Inches(1.0))

    # Arrow between cards (horizontal)
    if col < 3:
        add_flow_arrow(slide, x + Inches(3.0), y + Inches(1.2), Inches(0.25), color)

# Bottom summary
add_text(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.35),
         "최종 결과: 거래처(S001) → 매입상품(stock:0→50) → 판매(stock:50→47) | StockMovement: PURCHASE_IN(+50), SALE_OUT(-3)",
         11, ACCENT_TEAL, False)


# ════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "asis", "POSON_TOBE_매입_판매_플로우_분석.pptx")
prs.save(output_path)
print(f"PPTX saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
