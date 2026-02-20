"""설정 항목 전체 목록 PPTX 생성"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Color Palette ───
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1E, 0x29, 0x3B)
BG_ROW = RGBColor(0x16, 0x20, 0x33)
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)
ACCENT_GREEN = RGBColor(0x4A, 0xDE, 0x80)
ACCENT_ORANGE = RGBColor(0xFB, 0xBF, 0x24)
ACCENT_RED = RGBColor(0xF8, 0x71, 0x71)
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
ACCENT_TEAL = RGBColor(0x2D, 0xD4, 0xBF)
ACCENT_PINK = RGBColor(0xF4, 0x72, 0xB6)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
TEXT_LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
BORDER_COLOR = RGBColor(0x33, 0x41, 0x55)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=14, color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.font.name = "Malgun Gothic"
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=12, color=TEXT_WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, text_color, is_bold, sz) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz or font_size)
        p.font.color.rgb = text_color or color
        p.font.bold = is_bold
        p.font.name = "Malgun Gothic"
        p.space_after = Pt(1)
    return txBox


def add_badge(slide, left, top, text, color=ACCENT_BLUE, width=Inches(1.2)):
    shape = add_shape(slide, left, top, width, Inches(0.32), fill_color=color)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(10)
    shape.text_frame.paragraphs[0].font.color.rgb = TEXT_WHITE
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "Malgun Gothic"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


def add_stat_card(slide, x, y, number, label, color, w=Inches(1.9), h=Inches(1.3)):
    card = add_shape(slide, x, y, w, h, fill_color=BG_CARD, border_color=color)
    add_text(slide, x + Inches(0.1), y + Inches(0.15), w - Inches(0.2), Inches(0.6),
             str(number), 32, color, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), y + Inches(0.8), w - Inches(0.2), Inches(0.4),
             label, 11, TEXT_GRAY, False, PP_ALIGN.CENTER)
    return card


def add_table_row(slide, x, y, cols, widths, colors=None, bold=False, size=10):
    """간단한 텍스트 기반 테이블 행"""
    for i, (text, w) in enumerate(zip(cols, widths)):
        c = (colors[i] if colors else TEXT_LIGHT) if not bold else TEXT_WHITE
        add_text(slide, x, y, w, Inches(0.3), text, size, c, bold)
        x += w


def add_setting_items(slide, x, y, items, accent_color, max_items=20, col_width=Inches(5.8)):
    """설정 항목 리스트를 컴팩트하게 표시"""
    for i, (key, title, type_label) in enumerate(items[:max_items]):
        row_y = y + Inches(i * 0.26)
        type_color = {
            "toggle": ACCENT_GREEN,
            "number": ACCENT_BLUE,
            "text": ACCENT_ORANGE,
            "select": ACCENT_PURPLE,
            "date": ACCENT_TEAL,
            "data": TEXT_GRAY,
        }.get(type_label, TEXT_GRAY)
        add_text(slide, x, row_y, Inches(1.6), Inches(0.24), key, 8, TEXT_GRAY, False)
        add_text(slide, x + Inches(1.65), row_y, Inches(2.5), Inches(0.24), title, 9, TEXT_LIGHT, False)
        add_text(slide, x + Inches(4.3), row_y, Inches(0.8), Inches(0.24), type_label, 8, type_color, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_TEAL)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
         "POSON POS 시스템", 20, TEXT_GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.3), Inches(11), Inches(1.2),
         "설정 항목 전체 목록", 44, TEXT_WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.7), Inches(11), Inches(0.6),
         "공통 환경설정 (SettingsView) + 기기별 설정 (DevicesView)", 16, ACCENT_TEAL, False, PP_ALIGN.CENTER)

# 3 info cards
for i, (label, value) in enumerate([
    ("총 설정 항목", "319개"),
    ("DB 테이블", "SystemSetting + DeviceSetting"),
    ("분석일", "2026-02-19"),
]):
    x = Inches(2.5 + i * 3)
    card = add_shape(slide, x, Inches(5.0), Inches(2.6), Inches(1.1), fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text(slide, x + Inches(0.2), Inches(5.1), Inches(2.2), Inches(0.4), label, 11, TEXT_GRAY, False)
    add_text(slide, x + Inches(0.2), Inches(5.55), Inches(2.2), Inches(0.5), value, 14, TEXT_WHITE, True)


# ════════════════════════════════════════════════════════════════
# SLIDE 2: Overview Summary
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_TEAL)
add_text(slide, Inches(0.5), Inches(0.25), Inches(10), Inches(0.5), "전체 현황 요약", 28, TEXT_WHITE, True)

# Top stat cards - 큰 숫자
add_stat_card(slide, Inches(0.5), Inches(1.0), "319", "전체 항목", ACCENT_TEAL, Inches(2.0), Inches(1.4))
add_stat_card(slide, Inches(2.7), Inches(1.0), "102", "공통 환경설정", ACCENT_BLUE, Inches(2.0), Inches(1.4))
add_stat_card(slide, Inches(4.9), Inches(1.0), "217", "기기별 설정", ACCENT_GREEN, Inches(2.0), Inches(1.4))
add_stat_card(slide, Inches(7.1), Inches(1.0), "215", "토글 항목", ACCENT_ORANGE, Inches(2.0), Inches(1.4))
add_stat_card(slide, Inches(9.3), Inches(1.0), "63", "입력 항목", ACCENT_PURPLE, Inches(2.0), Inches(1.4))
add_stat_card(slide, Inches(11.5), Inches(1.0), "41", "선택 항목", ACCENT_PINK, Inches(2.0), Inches(1.4))

# 공통 환경설정 테이블
add_text(slide, Inches(0.5), Inches(2.7), Inches(5), Inches(0.4), "공통 환경설정 (SystemSetting)", 16, ACCENT_BLUE, True)
add_shape(slide, Inches(0.4), Inches(3.15), Inches(6.0), Inches(0.35), fill_color=ACCENT_BLUE)
add_text(slide, Inches(0.6), Inches(3.17), Inches(1.8), Inches(0.3), "탭", 10, TEXT_WHITE, True)
add_text(slide, Inches(2.4), Inches(3.17), Inches(1.5), Inches(0.3), "API 카테고리", 10, TEXT_WHITE, True)
add_text(slide, Inches(4.2), Inches(3.17), Inches(0.8), Inches(0.3), "항목 수", 10, TEXT_WHITE, True)
add_text(slide, Inches(5.1), Inches(3.17), Inches(1.2), Inches(0.3), "비중", 10, TEXT_WHITE, True)

common_tabs = [
    ("판매 운영", "SALE", "31", "30%", ACCENT_ORANGE),
    ("결제 정책", "PAYMENT", "21", "21%", ACCENT_RED),
    ("출력 설정", "PRINT", "9", "9%", ACCENT_BLUE),
    ("포인트/회원", "POINT", "27", "26%", ACCENT_GREEN),
    ("바코드/중량", "BARCODE", "5", "5%", ACCENT_TEAL),
    ("시스템", "SYSTEM", "9", "9%", ACCENT_PURPLE),
]
for j, (tab, api, cnt, pct, color) in enumerate(common_tabs):
    row_y = Inches(3.55 + j * 0.36)
    if j % 2 == 0:
        add_shape(slide, Inches(0.4), row_y, Inches(6.0), Inches(0.36), fill_color=BG_ROW)
    add_text(slide, Inches(0.6), row_y + Inches(0.03), Inches(1.8), Inches(0.3), tab, 11, TEXT_LIGHT, False)
    add_text(slide, Inches(2.4), row_y + Inches(0.03), Inches(1.5), Inches(0.3), api, 10, color, False)
    add_text(slide, Inches(4.2), row_y + Inches(0.03), Inches(0.8), Inches(0.3), cnt, 12, TEXT_WHITE, True)
    add_text(slide, Inches(5.1), row_y + Inches(0.03), Inches(1.2), Inches(0.3), pct, 10, TEXT_GRAY, False)

# 기기별 설정 테이블
add_text(slide, Inches(7.0), Inches(2.7), Inches(5), Inches(0.4), "기기별 설정 (DeviceSetting)", 16, ACCENT_GREEN, True)
add_shape(slide, Inches(6.9), Inches(3.15), Inches(6.0), Inches(0.35), fill_color=ACCENT_GREEN)
add_text(slide, Inches(7.1), Inches(3.17), Inches(1.5), Inches(0.3), "기기 유형", 10, TEXT_WHITE, True)
add_text(slide, Inches(8.6), Inches(3.17), Inches(0.8), Inches(0.3), "탭 수", 10, TEXT_WHITE, True)
add_text(slide, Inches(9.5), Inches(3.17), Inches(0.8), Inches(0.3), "항목 수", 10, TEXT_WHITE, True)
add_text(slide, Inches(10.4), Inches(3.17), Inches(2.4), Inches(0.3), "주요 탭", 10, TEXT_WHITE, True)

device_rows = [
    ("POS", "6", "153", "터미널(30) VAN(9) 인쇄(76) 판매(11) 정산(27) 영수증(10)", ACCENT_BLUE),
    ("KIOSK", "8", "86", "터미널(30) VAN(9) 현금(10) 봉투(6) 자동(8) 포인트(11) 인쇄(4) 기타(8)", ACCENT_GREEN),
    ("KITCHEN", "2", "37", "터미널(30) 주방메시지(7)", ACCENT_ORANGE),
]
for j, (dev, tabs, cnt, detail, color) in enumerate(device_rows):
    row_y = Inches(3.55 + j * 0.5)
    add_shape(slide, Inches(6.9), row_y, Inches(6.0), Inches(0.5), fill_color=BG_CARD, border_color=color)
    add_badge(slide, Inches(7.0), row_y + Inches(0.09), dev, color, Inches(0.9))
    add_text(slide, Inches(8.6), row_y + Inches(0.1), Inches(0.8), Inches(0.3), tabs, 14, TEXT_WHITE, True)
    add_text(slide, Inches(9.5), row_y + Inches(0.1), Inches(0.8), Inches(0.3), cnt, 14, color, True)
    add_text(slide, Inches(10.4), row_y + Inches(0.1), Inches(2.5), Inches(0.3), detail, 8, TEXT_GRAY, False)

# 아키텍처 요약
add_text(slide, Inches(0.5), Inches(5.8), Inches(12), Inches(0.4), "2계층 설정 아키텍처", 16, TEXT_WHITE, True)
arch_card = add_shape(slide, Inches(0.4), Inches(6.25), Inches(12.5), Inches(1.0), fill_color=BG_CARD, border_color=BORDER_COLOR)
arch_items = [
    ("공통 (SystemSetting)", "GET/PUT /api/v1/settings/:category", "매장 전체에서 공유하는 설정", ACCENT_BLUE),
    ("기기별 (DeviceSetting)", "GET/PUT /api/v1/devices/:id/settings/:category", "개별 기기에서만 적용되는 설정", ACCENT_GREEN),
]
for j, (name, api, desc, color) in enumerate(arch_items):
    ax = Inches(0.7 + j * 6.2)
    add_text(slide, ax, Inches(6.35), Inches(3), Inches(0.3), name, 13, color, True)
    add_text(slide, ax, Inches(6.65), Inches(5.5), Inches(0.25), api, 9, TEXT_GRAY, False)
    add_text(slide, ax, Inches(6.9), Inches(5.5), Inches(0.25), desc, 10, TEXT_LIGHT, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 3: 공통 - 판매 운영 (31항목)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_BLUE)
add_badge(slide, Inches(0.5), Inches(0.2), "공통", ACCENT_BLUE, Inches(0.7))
add_text(slide, Inches(1.3), Inches(0.15), Inches(8), Inches(0.5), "판매 운영 (SALE) — 31항목", 24, TEXT_WHITE, True)
add_text(slide, Inches(0.5), Inches(0.7), Inches(10), Inches(0.3), "API: GET/PUT /api/v1/settings/sale  |  키 프리픽스: sale.*", 11, TEXT_GRAY, False)

# 좌측: 영업관리 + 가격/금액 + 주방/테이블
add_text(slide, Inches(0.5), Inches(1.1), Inches(4), Inches(0.3), "영업 관리", 13, ACCENT_ORANGE, True)
mgmt_items = [
    ("openDay", "영업 시작일", "date"),
    ("finishDay", "마감일", "date"),
    ("receiptSeq", "전표 시퀀스", "number"),
    ("receiptNumber", "전표번호", "text"),
    ("startPrice", "시재금", "number"),
    ("beforTran", "이전 거래", "text"),
]
add_setting_items(slide, Inches(0.5), Inches(1.4), mgmt_items, ACCENT_ORANGE)

add_text(slide, Inches(0.5), Inches(3.0), Inches(4), Inches(0.3), "가격/금액", 13, ACCENT_ORANGE, True)
price_items = [
    ("maxPrice", "최대 결제금액 (9,999,990)", "number"),
    ("maxCashPrice", "최대 현금금액", "number"),
    ("delay", "지연 설정", "number"),
]
add_setting_items(slide, Inches(0.5), Inches(3.3), price_items, ACCENT_ORANGE)

add_text(slide, Inches(0.5), Inches(4.1), Inches(4), Inches(0.3), "주방/테이블", 13, ACCENT_ORANGE, True)
kitchen_items = [
    ("kitchenCallEnabled", "주방 호출", "toggle"),
    ("tableSelectEnabled", "테이블 선택", "toggle"),
    ("tableCount", "테이블 갯수", "number"),
]
add_setting_items(slide, Inches(0.5), Inches(4.4), kitchen_items, ACCENT_ORANGE)

# 우측: 판매 옵션 토글 13개 + 데이터 전용 6개
add_text(slide, Inches(6.5), Inches(1.1), Inches(6), Inches(0.3), "판매 옵션 (토글)", 13, ACCENT_GREEN, True)
sale_toggle_items = [
    ("priceEditable", "가격 수정 허용", "toggle"),
    ("productSound", "상품 스캔 소리", "toggle"),
    ("orderCallEnabled", "주문 호출", "toggle"),
    ("totalQtyShow", "총 수량 표시", "toggle"),
    ("grouping", "상품 그룹핑", "toggle"),
    ("boryuEnabled", "보류 기능", "toggle"),
    ("infoDeskEnabled", "안내데스크 모드", "toggle"),
    ("allFinish", "전체 마감", "toggle"),
    ("jobFinishCashdraw", "마감시 현금함", "toggle"),
    ("freeOpt", "무료 옵션", "toggle"),
    ("price11", "3단 가격 표시", "toggle"),
    ("engEnabled", "영어 모드", "toggle"),
    ("gridFix", "그리드 고정", "toggle"),
]
add_setting_items(slide, Inches(6.5), Inches(1.4), sale_toggle_items, ACCENT_GREEN)

add_text(slide, Inches(6.5), Inches(4.8), Inches(6), Inches(0.3), "데이터 전용 (UI 미노출)", 13, TEXT_GRAY, True)
data_items = [
    ("saleView", "판매 화면 유형", "data"),
    ("gridSaleEx", "그리드 확장", "data"),
    ("boryuTranOpt", "보류 전송 옵션", "data"),
    ("infoDeskViewAll", "안내데스크 전체보기", "data"),
    ("saleFinishOpt", "판매 마감 옵션", "data"),
    ("dayFinishMsgOpt", "일마감 메시지", "data"),
]
add_setting_items(slide, Inches(6.5), Inches(5.1), data_items, TEXT_GRAY)

# 범례
add_shape(slide, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.5), fill_color=BG_CARD, border_color=BORDER_COLOR)
legend = [
    ("toggle", ACCENT_GREEN), ("number", ACCENT_BLUE), ("text", ACCENT_ORANGE),
    ("select", ACCENT_PURPLE), ("date", ACCENT_TEAL), ("data", TEXT_GRAY),
]
for i, (lbl, c) in enumerate(legend):
    lx = Inches(0.7 + i * 2.0)
    add_badge(slide, lx, Inches(6.88), lbl, c, Inches(0.8))
    add_text(slide, lx + Inches(0.85), Inches(6.88), Inches(1), Inches(0.3), f"= {lbl}", 9, TEXT_GRAY, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 4: 공통 - 결제 정책 (21항목)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_BLUE)
add_badge(slide, Inches(0.5), Inches(0.2), "공통", ACCENT_BLUE, Inches(0.7))
add_text(slide, Inches(1.3), Inches(0.15), Inches(8), Inches(0.5), "결제 정책 (PAYMENT) — 21항목", 24, TEXT_WHITE, True)
add_text(slide, Inches(0.5), Inches(0.7), Inches(10), Inches(0.3), "API: GET/PUT /api/v1/settings/payment  |  키 프리픽스: payment.*", 11, TEXT_GRAY, False)

# 좌측: 카드/결제 정책 토글 9개
add_text(slide, Inches(0.5), Inches(1.1), Inches(5), Inches(0.3), "카드/결제 정책 (9개)", 13, ACCENT_RED, True)
card_toggles = [
    ("offCardCheck", "오프라인 카드 결제", "toggle"),
    ("offCardKeyUse", "오프라인 카드 키", "toggle"),
    ("handCardEnabled", "수기 카드 입력", "toggle"),
    ("cardTimerEnabled", "카드 타이머", "toggle"),
    ("cardWavOpt", "카드 결제 소리", "toggle"),
    ("cardView", "카드 정보 표시", "toggle"),
    ("eCardEnabled", "전자카드 사용", "toggle"),
    ("noCvmBillPrint", "비CVM 영수증", "toggle"),
    ("minCardPrice", "카드 최소금액", "number"),
]
add_setting_items(slide, Inches(0.5), Inches(1.4), card_toggles, ACCENT_RED)

# VAN/현금 + 상품권
add_text(slide, Inches(0.5), Inches(3.8), Inches(5), Inches(0.3), "VAN/현금/상품권 (4개)", 13, ACCENT_ORANGE, True)
van_items = [
    ("cashBackEnabled", "캐시백 사용", "toggle"),
    ("oCashScreen", "현금영수증 화면", "toggle"),
    ("giftInputEnabled", "상품권 입력", "toggle"),
    ("giftBillEtc", "상품권 기타 처리", "toggle"),
]
add_setting_items(slide, Inches(0.5), Inches(4.1), van_items, ACCENT_ORANGE)

# 환불
add_text(slide, Inches(0.5), Inches(5.15), Inches(5), Inches(0.3), "환불 (3개)", 13, ACCENT_PURPLE, True)
refund_items = [
    ("rePoint", "환불 포인트 재계산", "toggle"),
    ("reTax", "환불 세금 재계산", "toggle"),
    ("reCashBack", "환불 캐시백 재계산", "toggle"),
]
add_setting_items(slide, Inches(0.5), Inches(5.45), refund_items, ACCENT_PURPLE)

# 우측: 수수료 설정
add_text(slide, Inches(7.0), Inches(1.1), Inches(5), Inches(0.3), "수수료 설정 (5개)", 13, ACCENT_TEAL, True)
comm_items = [
    ("commCard", "카드 수수료 (%)", "number"),
    ("commPoint", "포인트 수수료 (%)", "number"),
    ("commCashBack", "캐시백 수수료 (%)", "number"),
    ("commCash", "현금 수수료 (%)", "number"),
    ("commCashRate", "현금 비율 (%)", "number"),
]
add_setting_items(slide, Inches(7.0), Inches(1.4), comm_items, ACCENT_TEAL)


# ════════════════════════════════════════════════════════════════
# SLIDE 5: 공통 - 출력/포인트/바코드/시스템
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_BLUE)
add_badge(slide, Inches(0.5), Inches(0.2), "공통", ACCENT_BLUE, Inches(0.7))
add_text(slide, Inches(1.3), Inches(0.15), Inches(10), Inches(0.5), "출력 / 포인트·회원 / 바코드·중량 / 시스템", 22, TEXT_WHITE, True)

# ─ 출력 설정 (9) ─
add_badge(slide, Inches(0.4), Inches(0.75), "PRINT 9개", ACCENT_BLUE, Inches(1.2))
card = add_shape(slide, Inches(0.4), Inches(1.1), Inches(6.2), Inches(2.7), fill_color=BG_CARD, border_color=ACCENT_BLUE)
print_items = [
    ("printVat", "부가세 인쇄", "toggle"),
    ("printBarcode", "바코드 인쇄", "toggle"),
    ("bottomPrint", "하단 문구 인쇄", "toggle"),
    ("pointBillPrint", "포인트 영수증", "toggle"),
    ("reTranBillPrint", "재거래 영수증", "toggle"),
    ("memberReceiptPrint", "회원 미수 인쇄", "toggle"),
    ("printerOffCheck", "프린터 오프 체크", "toggle"),
    ("slotAdd", "용지 슬롯 추가", "toggle"),
    ("cutPosition", "절단 위치 (기본/위/아래)", "select"),
]
add_setting_items(slide, Inches(0.6), Inches(1.2), print_items, ACCENT_BLUE)

# ─ 바코드/중량 (5) ─
add_badge(slide, Inches(6.9), Inches(0.75), "BARCODE 5개", ACCENT_TEAL, Inches(1.4))
card = add_shape(slide, Inches(6.9), Inches(1.1), Inches(6.0), Inches(1.6), fill_color=BG_CARD, border_color=ACCENT_TEAL)
barcode_items = [
    ("barCodeLen", "바코드 자동부여 숫자 (95)", "text"),
    ("scaleLen", "중량상품 코드 길이 (4)", "text"),
    ("scaleStartChar", "중량상품 시작문자 (28)", "text"),
    ("scalePriceCut", "중량 가격 절사", "select"),
    ("scale18YN", "18자리 중량 바코드", "toggle"),
]
add_setting_items(slide, Inches(7.1), Inches(1.2), barcode_items, ACCENT_TEAL)

# ─ 시스템 (9) ─
add_badge(slide, Inches(6.9), Inches(2.85), "SYSTEM 9개", ACCENT_PURPLE, Inches(1.4))
card = add_shape(slide, Inches(6.9), Inches(3.2), Inches(6.0), Inches(2.6), fill_color=BG_CARD, border_color=ACCENT_PURPLE)
sys_items = [
    ("errorLog", "오류 로그 기록", "toggle"),
    ("mdbCompact", "DB 자동 컴팩트", "toggle"),
    ("masterDownEnabled", "마스터 자동 다운", "toggle"),
    ("masterDownWeek", "마스터 다운 주기 (주)", "number"),
    ("newItemUpdate", "신규상품 업데이트", "toggle"),
    ("scanRealCheck", "실시간 스캔 체크", "toggle"),
    ("logoMinimize", "로고 최소화", "toggle"),
    ("screenHide", "화면 숨김", "toggle"),
    ("backupPath", "백업 경로", "text"),
]
add_setting_items(slide, Inches(7.1), Inches(3.3), sys_items, ACCENT_PURPLE)

# ─ 포인트/회원 (27) - 컴팩트 ─
add_badge(slide, Inches(0.4), Inches(3.95), "POINT 27개", ACCENT_GREEN, Inches(1.3))
card = add_shape(slide, Inches(0.4), Inches(4.3), Inches(6.2), Inches(3.0), fill_color=BG_CARD, border_color=ACCENT_GREEN)

add_text(slide, Inches(0.6), Inches(4.35), Inches(3), Inches(0.25), "포인트 기본 (7개)", 10, ACCENT_GREEN, True)
point_basic = [
    ("salePoint", "판매 포인트", "toggle"),
    ("weightPoint", "중량 상품 포인트", "select"),
    ("memberAddScreen", "회원 추가 화면", "toggle"),
    ("gradeMemo / noBillMessage / noBillSound / noBillCusPoint", "", "toggle"),
]
for i, (k, t, tp) in enumerate(point_basic):
    add_text(slide, Inches(0.6), Inches(4.6 + i * 0.22), Inches(5.8), Inches(0.2), f"{k}  {t}", 8, TEXT_LIGHT if t else TEXT_GRAY, False)

add_text(slide, Inches(0.6), Inches(5.5), Inches(3), Inches(0.25), "고객 UI - 키오스크 (20개)", 10, ACCENT_GREEN, True)
self_ui_summary = [
    "입력: selfCusTopMsg, selfCusBTMsg1, selfCusBTMsg2",
    "선택: selfCusSelect, selfReader, selfStartHotKey, selfPriceType, selfCusAddEtc",
    "토글: selfSoundGuide, selfCusNum4, selfNoCustomer, selfCusAddUse,",
    "      selfTouchSoundYN, selfMainPage, selfBTInit, selfOneCancel,",
    "      selfZHotKey, selfCountYN, selfPriceUse",
]
for i, line in enumerate(self_ui_summary):
    color = ACCENT_ORANGE if "입력" in line else (ACCENT_PURPLE if "선택" in line else TEXT_GRAY)
    add_text(slide, Inches(0.6), Inches(5.75 + i * 0.22), Inches(5.8), Inches(0.2), line, 8, color, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 6: 기기 - 터미널 HW + VAN (공통)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_GREEN)
add_badge(slide, Inches(0.5), Inches(0.2), "기기", ACCENT_GREEN, Inches(0.7))
add_text(slide, Inches(1.3), Inches(0.15), Inches(10), Inches(0.5), "터미널 HW (30항목) + VAN 결제 (9항목)", 24, TEXT_WHITE, True)
add_text(slide, Inches(0.5), Inches(0.7), Inches(10), Inches(0.3), "POS·KIOSK·KITCHEN 공통  |  API: /api/v1/devices/:id/settings/terminal, van", 11, TEXT_GRAY, False)

# 터미널 - 좌측
add_text(slide, Inches(0.5), Inches(1.1), Inches(5), Inches(0.3), "터미널 HW (TERMINAL) — 30항목", 14, ACCENT_BLUE, True)
card = add_shape(slide, Inches(0.4), Inches(1.4), Inches(6.2), Inches(5.8), fill_color=BG_CARD, border_color=ACCENT_BLUE)

terminal_groups = [
    ("기본정보", [("terminalType", "터미널 유형"), ("posNo", "POS 번호"), ("adminPosNo", "관리자 POS번호")]),
    ("주변기기", [("cashDraw", "캐시 드로워 [T]"), ("touch", "터치스크린 [T]"), ("dual", "듀얼 모니터 [T]"), ("dualType", "듀얼 유형")]),
    ("프린터", [("printer", "프린터 종류"), ("printerPort", "프린터 포트"), ("printerBps", "전송속도")]),
    ("스캐너", [("scanName", "스캐너 종류"), ("scanPort", "스캐너 포트"), ("handScanPort", "핸드스캐너 포트")]),
    ("MSR", [("msrPort", "MSR 포트"), ("msrBps", "MSR 전송속도")]),
    ("CDP", [("cdpName", "CDP 이름"), ("cdpPort", "CDP 포트"), ("cdpLine", "CDP 라인"), ("cdpType", "CDP 유형"), ("cdpBps", "CDP 속도"), ("cdpCashYN", "CDP 현금표시 [T]")]),
    ("기타", [("coinName", "코인기"), ("coinPort", "코인기 포트"), ("moniter", "모니터"), ("sMoniter", "보조모니터"), ("printerCatUse", "CAT프린터 [T]"), ("catPort", "CAT포트"), ("catBps", "CAT속도"), ("cboScalePort", "저울 포트"), ("supyoPort", "수표 포트")]),
]
row_y = Inches(1.5)
for group_name, items in terminal_groups:
    add_text(slide, Inches(0.6), row_y, Inches(2), Inches(0.22), group_name, 9, ACCENT_BLUE, True)
    for k, v in items:
        row_y += Inches(0.22)
        marker = ACCENT_GREEN if "[T]" in v else TEXT_GRAY
        add_text(slide, Inches(0.8), row_y, Inches(2.2), Inches(0.2), k, 8, TEXT_GRAY, False)
        add_text(slide, Inches(3.1), row_y, Inches(3), Inches(0.2), v.replace(" [T]", ""), 8, TEXT_LIGHT, False)
        if "[T]" in v:
            add_text(slide, Inches(5.5), row_y, Inches(0.8), Inches(0.2), "toggle", 7, ACCENT_GREEN, False)
    row_y += Inches(0.28)

# VAN - 우측
add_text(slide, Inches(7.0), Inches(1.1), Inches(5), Inches(0.3), "VAN 결제 (VAN) — 9항목", 14, ACCENT_RED, True)
card = add_shape(slide, Inches(6.9), Inches(1.4), Inches(6.0), Inches(3.0), fill_color=BG_CARD, border_color=ACCENT_RED)

van_items = [
    ("vanSelect", "VAN사 선택 (NICE/KSNET/KICC/KIS...)", "select"),
    ("singPadPort", "서명패드 포트", "select"),
    ("logDelete", "로그 삭제", "toggle"),
    ("vanIp", "VAN IP", "text"),
    ("vanPort", "VAN 포트", "text"),
    ("danmalNo", "단말기 번호 (TID)", "text"),
    ("snumber", "시리얼 번호", "text"),
    ("singPadBps", "서명패드 전송속도", "select"),
    ("usbYN", "USB 사용", "toggle"),
]
add_setting_items(slide, Inches(7.1), Inches(1.55), van_items, ACCENT_RED)


# ════════════════════════════════════════════════════════════════
# SLIDE 7: POS 인쇄/출력 (76항목)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_ORANGE)
add_badge(slide, Inches(0.5), Inches(0.2), "POS", ACCENT_BLUE, Inches(0.7))
add_text(slide, Inches(1.3), Inches(0.15), Inches(10), Inches(0.5), "인쇄/출력 (POS_PRINT) — 76항목 (전체 토글)", 22, TEXT_WHITE, True)
add_text(slide, Inches(0.5), Inches(0.7), Inches(10), Inches(0.3), "API: /api/v1/devices/:id/settings/pos_print  |  ASIS: S_Config 79개 항목 매핑", 11, TEXT_GRAY, False)

# 3열로 분할
pos_print_all = [
    "1.정산 분류별 매출 출력", "2.정산 금액단위별 수량", "3.신상품 등록 가능",
    "4.상품정보 1줄 표시", "5.환전 비밀번호", "6.0원 상품 입력",
    "7.할인상품 적립 안함", "8.중량상품 적립 안함", "9.전체취소 비활성",
    "10.전표번호 미출력", "11.배달기능 비활성", "12.보류 영수증 출력",
    "13.카드 시 금고 안열기", "14.취소 내역 출력", "15.외상 매출 불가",
    "16.영수증 바코드", "17.영수증 부가세", "18.공병합계 표시",
    "19.캐쉬백 사용", "20.매입가 미표시", "21.바코드별 상품 묶기",
    "22.반품 비밀번호 미사용", "23.회원 구매금액 숨김", "24.카드 보관용 상세인쇄",
    "25.수표 응답 인쇄", "26.배달 포인트 적립",
    "27.할인상품 캐쉬백 안함", "28.받은돈 무조건 입력", "29.보관용 바코드 안함",
    "30.전자서명 보관전표 안함", "31.배달 보관전표 2장", "32.재발행 바코드 안함",
    "33.전표번호 숨기기", "34.내방배달 보관전표", "35.배달미입금 마감불가",
    "36.마감시 외상입금 출력", "37.미수고객 경고창", "38.보류시 교대 불가",
    "39.배달미처리 교대 불가", "40.상품순번 출력", "41.현금전용 반품 불가",
    "42.스캔시 보류 메시지",
    "44.18행 중량바코드", "45.보류 영수증 상세", "46.현금영수증 자진발급",
    "47.매입>판매 경고", "48.판매메세지 숨김", "49.중지상품 변경사용",
    "50.배달 약식 출력", "51.미수금 출력 안함", "52.배달미입금 반품 가능",
    "53.배달판매 구분", "54.가격비교 사용",
    "57.부류/저울 판매가 입력", "58.배달 비고란 출력", "59.바코드 카드결제 화면",
    "60.잔액 카드결제 화면", "61.단축키 고정판매가", "62.상품권 현금영수증 포함",
    "63.상품권 포인트 포함", "64.개점/교대 재정정보", "65.배달 순번 출력",
    "66.배달 영수증 관리", "67.기사용 영수증 출력", "68.벌크상품 합계 맞춤",
    "69.캐쉬백 QR 인쇄", "70.전화번호 마스킹", "71.단축키 마포이름",
    "72.e-쿠폰 영수증 안함", "73.배달 개별상품 보관용", "74.현금영수증 고객표시",
    "75.배달순번 회원영수증", "76.회원명 마스킹", "77.카드반품 돈통 안열기",
    "78.카드+카드 결제", "79.배달 보관 일반설정",
]

# 3 columns
col_size = 26
for col in range(3):
    cx = Inches(0.3 + col * 4.3)
    card = add_shape(slide, cx, Inches(1.05), Inches(4.1), Inches(6.2), fill_color=BG_CARD, border_color=BORDER_COLOR)
    start = col * col_size
    end = min(start + col_size, len(pos_print_all))
    for j, item in enumerate(pos_print_all[start:end]):
        color = ACCENT_GREEN if j < 13 else (ACCENT_ORANGE if j < 26 else TEXT_LIGHT)
        add_text(slide, cx + Inches(0.1), Inches(1.15 + j * 0.235), Inches(3.9), Inches(0.22), item, 9, color if col == 0 else TEXT_LIGHT, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 8: POS 판매/정산/영수증
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_ORANGE)
add_badge(slide, Inches(0.5), Inches(0.2), "POS", ACCENT_BLUE, Inches(0.7))
add_text(slide, Inches(1.3), Inches(0.15), Inches(10), Inches(0.5), "판매설정(11) + 정산/마감(27) + 영수증(10)", 22, TEXT_WHITE, True)

# 판매설정 (11) - 좌측 상단
add_badge(slide, Inches(0.4), Inches(0.75), "POS_SALE 11개", ACCENT_ORANGE, Inches(1.5))
card = add_shape(slide, Inches(0.4), Inches(1.1), Inches(4.0), Inches(3.4), fill_color=BG_CARD, border_color=ACCENT_ORANGE)
pos_sale_items = [
    ("receiptDiscountDisplay", "영수증 할인 표시", "select"),
    ("totalRounding", "합계 반올림", "select"),
    ("scaleRounding", "저울 반올림", "select"),
    ("categoryPrintType", "분류 출력 유형", "select"),
    ("deliveryPointRate", "배달 포인트 비율 (%)", "number"),
    ("creditPointRate", "외상 포인트 비율 (%)", "number"),
    ("minReceiptAmount", "최소 영수증 금액", "number"),
    ("minPointAmount", "최소 포인트 금액", "number"),
    ("cardNoSignAmount", "카드 무서명 금액 (50,000)", "number"),
    ("bcPartner", "BC 파트너스", "toggle"),
    ("creditMemoUse", "외상결제 메모", "toggle"),
]
add_setting_items(slide, Inches(0.6), Inches(1.2), pos_sale_items, ACCENT_ORANGE)

# 영수증 (10) - 좌측 하단
add_badge(slide, Inches(0.4), Inches(4.65), "POS_RECEIPT 10개", ACCENT_PINK, Inches(1.6))
card = add_shape(slide, Inches(0.4), Inches(5.0), Inches(4.0), Inches(2.3), fill_color=BG_CARD, border_color=ACCENT_PINK)
receipt_items = [
    "배달시 출력 안함", "외상시 출력 안함", "포인트 사용 출력 안함",
    "상품권 사용 출력 안함", "e쿠폰 사용 출력 안함", "QRMS 사용 출력 안함",
    "카드 일시불 출력 안함", "카드 할부 출력 안함", "현금영수증 출력 안함",
    "전표반품 출력 안함",
]
add_text(slide, Inches(0.6), Inches(5.05), Inches(3.5), Inches(0.25), "결제 유형별 영수증 출력 억제 설정 (모두 toggle)", 9, TEXT_GRAY, False)
for i, item in enumerate(receipt_items):
    col = i // 5
    row = i % 5
    add_text(slide, Inches(0.6 + col * 2.0), Inches(5.3 + row * 0.22), Inches(1.9), Inches(0.2), f"{'OFF' if i < 5 else 'OFF'} {item}", 8, ACCENT_PINK, False)

# 정산/마감 (27) - 우측
add_badge(slide, Inches(4.7), Inches(0.75), "POS_SETTLE 27개", ACCENT_PURPLE, Inches(1.6))

# 담당자 정산서 (13)
add_text(slide, Inches(4.7), Inches(1.15), Inches(4), Inches(0.25), "담당자 정산서 (13개)", 11, ACCENT_PURPLE, True)
card = add_shape(slide, Inches(4.7), Inches(1.4), Inches(4.0), Inches(2.0), fill_color=BG_CARD, border_color=ACCENT_PURPLE)
staff_items = [
    "공병 매출", "과/면세 내역", "할인 내역", "취소 내역", "신용카드 내역",
    "현금영수증 내역", "경품지급 내역", "매입사별 카드매출", "외상 입금 내역",
    "전배달 입금 내역", "기타 할인 내역", "배달건수 출력", "정산서 출력",
]
for i, item in enumerate(staff_items):
    col = i // 7
    row = i % 7
    add_text(slide, Inches(4.9 + col * 2.0), Inches(1.5 + row * 0.24), Inches(1.9), Inches(0.2), item, 9, TEXT_LIGHT, False)

# 마감 정산서 (10)
add_text(slide, Inches(4.7), Inches(3.5), Inches(4), Inches(0.25), "마감 정산서 (10개)", 11, ACCENT_PURPLE, True)
card = add_shape(slide, Inches(4.7), Inches(3.75), Inches(4.0), Inches(1.6), fill_color=BG_CARD, border_color=ACCENT_PURPLE)
close_items = [
    "공병 매출", "과/면세 내역", "할인 내역", "취소 내역", "신용카드 내역",
    "현금영수증 내역", "경품지급 내역", "매입사별 카드매출", "배달건수 출력", "정산서 출력",
]
for i, item in enumerate(close_items):
    col = i // 5
    row = i % 5
    add_text(slide, Inches(4.9 + col * 2.0), Inches(3.85 + row * 0.24), Inches(1.9), Inches(0.2), item, 9, TEXT_LIGHT, False)

# 영수증 메시지 (4)
add_text(slide, Inches(4.7), Inches(5.45), Inches(4), Inches(0.25), "영수증 메시지 (4개)", 11, ACCENT_PURPLE, True)
card = add_shape(slide, Inches(4.7), Inches(5.7), Inches(4.0), Inches(1.3), fill_color=BG_CARD, border_color=ACCENT_PURPLE)
msg_items = [
    ("receiptTopMsg", "영수증 상단 메시지", "text"),
    ("receiptBottomMsg1", "영수증 하단 메시지 1", "text"),
    ("receiptBottomMsg2", "영수증 하단 메시지 2", "text"),
    ("receiptBottomMsg3", "영수증 하단 메시지 3", "text"),
]
add_setting_items(slide, Inches(4.9), Inches(5.8), msg_items, ACCENT_PURPLE)

# POS 합계 카드 - 우측
add_shape(slide, Inches(9.0), Inches(0.75), Inches(4.0), Inches(6.25), fill_color=BG_CARD, border_color=ACCENT_BLUE)
add_text(slide, Inches(9.2), Inches(0.85), Inches(3.5), Inches(0.4), "POS 기기설정 합계", 16, ACCENT_BLUE, True)

pos_summary = [
    ("터미널 HW", "30", ACCENT_BLUE),
    ("VAN 결제", "9", ACCENT_RED),
    ("인쇄/출력", "76", ACCENT_ORANGE),
    ("판매설정", "11", ACCENT_ORANGE),
    ("정산/마감", "27", ACCENT_PURPLE),
    ("영수증", "10", ACCENT_PINK),
]
for i, (name, cnt, color) in enumerate(pos_summary):
    ry = Inches(1.4 + i * 0.55)
    add_text(slide, Inches(9.4), ry, Inches(2), Inches(0.3), name, 13, TEXT_LIGHT, False)
    add_text(slide, Inches(11.5), ry, Inches(1.2), Inches(0.3), cnt, 18, color, True)
    add_shape(slide, Inches(9.3), ry + Inches(0.35), Inches(3.5), Inches(0.02), fill_color=BORDER_COLOR)

add_text(slide, Inches(9.4), Inches(4.8), Inches(2), Inches(0.4), "합계", 16, TEXT_WHITE, True)
add_text(slide, Inches(11.2), Inches(4.7), Inches(1.5), Inches(0.5), "153", 28, ACCENT_BLUE, True)
add_text(slide, Inches(9.4), Inches(5.2), Inches(3.5), Inches(0.3), "6개 탭 (POS 전용)", 11, TEXT_GRAY, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 9: KIOSK (86항목)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_GREEN)
add_badge(slide, Inches(0.5), Inches(0.2), "KIOSK", ACCENT_GREEN, Inches(0.8))
add_text(slide, Inches(1.4), Inches(0.15), Inches(10), Inches(0.5), "키오스크 기기설정 — 6개 전용 탭 (86항목)", 22, TEXT_WHITE, True)
add_text(slide, Inches(0.5), Inches(0.7), Inches(10), Inches(0.3), "API: /api/v1/devices/:id/settings/{self_cash|self_bag|self_auto|self_point|self_print|self_etc}", 10, TEXT_GRAY, False)

# 현금 결제 (10)
add_badge(slide, Inches(0.4), Inches(1.05), "SELF_CASH 10개", ACCENT_RED, Inches(1.5))
card = add_shape(slide, Inches(0.4), Inches(1.4), Inches(4.0), Inches(2.8), fill_color=BG_CARD, border_color=ACCENT_RED)
self_cash_items = [
    ("selfCash", "현금 결제 사용", "toggle"),
    ("selfCashPort", "현금기 포트", "select"),
    ("selfCashSleep", "현금기 대기", "select"),
    ("selfCashPhonNum", "전화번호", "text"),
    ("selfCashGubun", "현금 구분", "select"),
    ("selfNoHyunYoung", "실결제금액 숨김", "toggle"),
    ("selfOneHPUse", "1만원권 사용", "toggle"),
    ("self50HPUse", "5만원권 사용", "toggle"),
    ("self10000Use", "만원권 사용", "toggle"),
    ("selfNoCardUse", "카드 결제 비활성", "toggle"),
]
add_setting_items(slide, Inches(0.6), Inches(1.5), self_cash_items, ACCENT_RED)

# 봉투/저울 (6)
add_badge(slide, Inches(4.6), Inches(1.05), "SELF_BAG 6개", ACCENT_TEAL, Inches(1.4))
card = add_shape(slide, Inches(4.6), Inches(1.4), Inches(4.0), Inches(1.8), fill_color=BG_CARD, border_color=ACCENT_TEAL)
self_bag_items = [
    ("selfBagPort", "봉투기 포트", "select"),
    ("selfStartBag", "시작시 봉투", "toggle"),
    ("selfMBagSell", "복수 봉투 판매", "toggle"),
    ("selfLastBag", "마지막 봉투", "toggle"),
    ("selfScalePort", "저울 포트", "select"),
    ("selfScaleLimitG", "저울 한계 중량 (g)", "number"),
]
add_setting_items(slide, Inches(4.8), Inches(1.5), self_bag_items, ACCENT_TEAL)

# 자동 운영 (8)
add_badge(slide, Inches(8.8), Inches(1.05), "SELF_AUTO 8개", ACCENT_PURPLE, Inches(1.5))
card = add_shape(slide, Inches(8.8), Inches(1.4), Inches(4.1), Inches(2.3), fill_color=BG_CARD, border_color=ACCENT_PURPLE)
self_auto_items = [
    ("autoOpenYN", "자동 개점", "toggle"),
    ("autoFinishYN", "자동 마감", "toggle"),
    ("autoDay", "자동 운영 요일", "text"),
    ("autoAP", "AM/PM 구분", "select"),
    ("autoHH", "시간", "select"),
    ("autoMM", "분", "select"),
    ("autoID", "자동 ID", "text"),
    ("autoPass", "자동 비밀번호", "text"),
]
add_setting_items(slide, Inches(9.0), Inches(1.5), self_auto_items, ACCENT_PURPLE)

# 포인트/알림 (11)
add_badge(slide, Inches(0.4), Inches(4.35), "SELF_POINT 11개", ACCENT_ORANGE, Inches(1.5))
card = add_shape(slide, Inches(0.4), Inches(4.7), Inches(4.0), Inches(2.6), fill_color=BG_CARD, border_color=ACCENT_ORANGE)
self_point_items = [
    ("selfNoAutoPoint", "자동포인트 비활성", "toggle"),
    ("selfPointZero", "포인트 초기화", "toggle"),
    ("selfPointHidden", "포인트 숨김", "toggle"),
    ("selfPointSMSUse", "포인트 SMS", "toggle"),
    ("selfUserCall", "직원 호출", "toggle"),
    ("selfSMSAdmin", "관리자 SMS", "toggle"),
    ("selfKakao", "카카오 알림", "toggle"),
    ("selfZero", "제로 설정", "toggle"),
    ("selfCusAlarmUse", "고객 알람", "toggle"),
    ("selfCusAlarmTime", "고객 알람 시간", "select"),
    ("selfSNSGubun", "SNS 구분", "select"),
]
add_setting_items(slide, Inches(0.6), Inches(4.8), self_point_items, ACCENT_ORANGE)

# 인쇄/출력 (4)
add_badge(slide, Inches(4.6), Inches(3.35), "SELF_PRINT 4개", ACCENT_BLUE, Inches(1.4))
card = add_shape(slide, Inches(4.6), Inches(3.7), Inches(4.0), Inches(1.3), fill_color=BG_CARD, border_color=ACCENT_BLUE)
self_print_items = [
    ("selfAutoPrint", "자동 출력", "toggle"),
    ("selfStoPrint", "출력 중지", "toggle"),
    ("selfPrintAddress", "주소 출력", "toggle"),
    ("selfPrintPhon", "전화번호 출력", "toggle"),
]
add_setting_items(slide, Inches(4.8), Inches(3.8), self_print_items, ACCENT_BLUE)

# 기타 (8)
add_badge(slide, Inches(4.6), Inches(5.1), "SELF_ETC 8개", ACCENT_PINK, Inches(1.3))
card = add_shape(slide, Inches(4.6), Inches(5.45), Inches(4.0), Inches(1.85), fill_color=BG_CARD, border_color=ACCENT_PINK)
self_etc_items = [
    ("selfJPYN", "일본 모드 (동전교환)", "toggle"),
    ("selfBagJPPort", "일본 봉투기 포트", "select"),
    ("selfNoAutoGoods", "자동상품 비활성", "toggle"),
    ("selfAppCard", "앱카드", "toggle"),
    ("selfApple", "애플페이", "toggle"),
    ("selfCamUse", "카메라", "toggle"),
    ("selfICSiren", "IC 사이렌", "toggle"),
]
add_setting_items(slide, Inches(4.8), Inches(5.55), self_etc_items, ACCENT_PINK)

# KIOSK 합계
add_shape(slide, Inches(8.8), Inches(3.9), Inches(4.1), Inches(3.4), fill_color=BG_CARD, border_color=ACCENT_GREEN)
add_text(slide, Inches(9.0), Inches(4.0), Inches(3.5), Inches(0.3), "KIOSK 기기설정 합계", 14, ACCENT_GREEN, True)
kiosk_summary = [
    ("터미널 HW", "30", TEXT_LIGHT),
    ("VAN 결제", "9", TEXT_LIGHT),
    ("현금 결제", "10", ACCENT_RED),
    ("봉투/저울", "6", ACCENT_TEAL),
    ("자동 운영", "8", ACCENT_PURPLE),
    ("포인트/알림", "11", ACCENT_ORANGE),
    ("인쇄/출력", "4", ACCENT_BLUE),
    ("기타", "8", ACCENT_PINK),
]
for i, (name, cnt, color) in enumerate(kiosk_summary):
    ry = Inches(4.4 + i * 0.32)
    add_text(slide, Inches(9.2), ry, Inches(2), Inches(0.28), name, 10, TEXT_LIGHT, False)
    add_text(slide, Inches(11.5), ry, Inches(1), Inches(0.28), cnt, 12, color, True)

add_shape(slide, Inches(9.0), Inches(6.95), Inches(3.7), Inches(0.02), fill_color=ACCENT_GREEN)
add_text(slide, Inches(9.2), Inches(7.0), Inches(1.5), Inches(0.3), "합계", 13, TEXT_WHITE, True)
add_text(slide, Inches(11.3), Inches(6.95), Inches(1.5), Inches(0.35), "86", 20, ACCENT_GREEN, True)


# ════════════════════════════════════════════════════════════════
# SLIDE 10: KITCHEN + 전체 통계
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_TEAL)
add_text(slide, Inches(0.5), Inches(0.25), Inches(10), Inches(0.5), "KITCHEN 설정 + 전체 통계", 28, TEXT_WHITE, True)

# Kitchen (7)
add_badge(slide, Inches(0.4), Inches(1.0), "KITCHEN 7개", ACCENT_ORANGE, Inches(1.3))
card = add_shape(slide, Inches(0.4), Inches(1.35), Inches(5.0), Inches(2.2), fill_color=BG_CARD, border_color=ACCENT_ORANGE)
kitchen_items = [
    ("kitchenPrint", "주방 프린터 종류", "select"),
    ("kitchenPrintPort", "주방 프린터 포트", "select"),
    ("kitchenPrintBps", "주방 프린터 속도", "select"),
    ("kitchenFontsize", "폰트 크기", "select"),
    ("kitchenMsg1", "주방 메시지 1", "text"),
    ("kitchenMsg2", "주방 메시지 2", "text"),
    ("kitchenMsg4", "주방 메시지 4", "text"),
]
add_setting_items(slide, Inches(0.6), Inches(1.45), kitchen_items, ACCENT_ORANGE)

add_text(slide, Inches(0.6), Inches(3.0), Inches(4.5), Inches(0.3),
         "KITCHEN = 터미널 HW(30) + 주방 메시지(7) = 37항목", 10, TEXT_GRAY, False)

# 전체 통계 테이블
add_text(slide, Inches(0.4), Inches(3.6), Inches(5), Inches(0.4), "전체 항목 유형별 통계", 18, TEXT_WHITE, True)
card = add_shape(slide, Inches(0.4), Inches(4.0), Inches(5.0), Inches(3.2), fill_color=BG_CARD, border_color=BORDER_COLOR)

# Header
add_shape(slide, Inches(0.4), Inches(4.0), Inches(5.0), Inches(0.4), fill_color=ACCENT_TEAL)
add_text(slide, Inches(0.6), Inches(4.05), Inches(1.8), Inches(0.3), "분류", 11, TEXT_WHITE, True)
add_text(slide, Inches(2.4), Inches(4.05), Inches(0.7), Inches(0.3), "토글", 11, TEXT_WHITE, True)
add_text(slide, Inches(3.1), Inches(4.05), Inches(0.7), Inches(0.3), "입력", 11, TEXT_WHITE, True)
add_text(slide, Inches(3.8), Inches(4.05), Inches(0.7), Inches(0.3), "선택", 11, TEXT_WHITE, True)
add_text(slide, Inches(4.5), Inches(4.05), Inches(0.8), Inches(0.3), "합계", 11, TEXT_WHITE, True)

stats_data = [
    ("공통 환경설정", "68", "22", "12", "102", ACCENT_BLUE),
    ("POS 전용", "111", "13", "9", "133", ACCENT_ORANGE),
    ("KIOSK 전용", "31", "9", "16", "56", ACCENT_GREEN),
    ("공통 터미널+VAN", "5", "19", "15", "39", TEXT_LIGHT),
    ("KITCHEN 전용", "0", "3", "4", "7", ACCENT_RED),
]
for j, (label, t, i, s, total, color) in enumerate(stats_data):
    ry = Inches(4.45 + j * 0.4)
    if j % 2 == 0:
        add_shape(slide, Inches(0.4), ry, Inches(5.0), Inches(0.4), fill_color=BG_ROW)
    add_text(slide, Inches(0.6), ry + Inches(0.05), Inches(1.8), Inches(0.3), label, 10, color, False)
    add_text(slide, Inches(2.4), ry + Inches(0.05), Inches(0.7), Inches(0.3), t, 11, ACCENT_GREEN, False)
    add_text(slide, Inches(3.1), ry + Inches(0.05), Inches(0.7), Inches(0.3), i, 11, ACCENT_ORANGE, False)
    add_text(slide, Inches(3.8), ry + Inches(0.05), Inches(0.7), Inches(0.3), s, 11, ACCENT_PURPLE, False)
    add_text(slide, Inches(4.5), ry + Inches(0.05), Inches(0.8), Inches(0.3), total, 12, TEXT_WHITE, True)

# Total row
add_shape(slide, Inches(0.4), Inches(6.5), Inches(5.0), Inches(0.4), fill_color=ACCENT_TEAL)
add_text(slide, Inches(0.6), Inches(6.55), Inches(1.8), Inches(0.3), "전체 합계", 11, TEXT_WHITE, True)
add_text(slide, Inches(2.4), Inches(6.55), Inches(0.7), Inches(0.3), "215", 12, TEXT_WHITE, True)
add_text(slide, Inches(3.1), Inches(6.55), Inches(0.7), Inches(0.3), "63", 12, TEXT_WHITE, True)
add_text(slide, Inches(3.8), Inches(6.55), Inches(0.7), Inches(0.3), "41", 12, TEXT_WHITE, True)
add_text(slide, Inches(4.5), Inches(6.55), Inches(0.8), Inches(0.3), "319", 14, TEXT_WHITE, True)

# 우측: 큰 원형 요약 카드
summary_cards = [
    ("319", "전체 설정 항목", ACCENT_TEAL),
    ("102", "공통 환경설정", ACCENT_BLUE),
    ("217", "기기별 설정", ACCENT_GREEN),
]
for i, (num, label, color) in enumerate(summary_cards):
    cx = Inches(6.0 + i * 2.5)
    add_stat_card(slide, cx, Inches(1.0), num, label, color, Inches(2.2), Inches(1.5))

# API 요약
add_text(slide, Inches(6.0), Inches(2.8), Inches(7), Inches(0.4), "API 엔드포인트 요약", 16, TEXT_WHITE, True)
card = add_shape(slide, Inches(5.9), Inches(3.2), Inches(7.0), Inches(4.0), fill_color=BG_CARD, border_color=BORDER_COLOR)

api_rows = [
    ("공통", "GET /api/v1/settings/:category", "SALE, PAYMENT, PRINT, POINT, BARCODE, SYSTEM", ACCENT_BLUE),
    ("공통", "PUT /api/v1/settings/:category", "동일 카테고리 (key-value 일괄 저장)", ACCENT_BLUE),
    ("기기", "GET /api/v1/devices", "기기 목록 조회 (POS/KIOSK/KITCHEN)", ACCENT_GREEN),
    ("기기", "POST /api/v1/devices", "신규 기기 등록", ACCENT_GREEN),
    ("기기", "GET /api/v1/devices/:id/settings/:cat", "TERMINAL, VAN, POS_PRINT, POS_SALE...", ACCENT_GREEN),
    ("기기", "PUT /api/v1/devices/:id/settings/:cat", "SELF_CASH, SELF_BAG, SELF_AUTO...", ACCENT_GREEN),
]
for j, (scope, endpoint, desc, color) in enumerate(api_rows):
    ry = Inches(3.4 + j * 0.45)
    add_badge(slide, Inches(6.1), ry + Inches(0.05), scope, color, Inches(0.6))
    add_text(slide, Inches(6.8), ry + Inches(0.05), Inches(3.5), Inches(0.3), endpoint, 10, TEXT_WHITE, False)
    add_text(slide, Inches(6.1), ry + Inches(0.3), Inches(6.5), Inches(0.2), desc, 8, TEXT_GRAY, False)

# DB 테이블
add_text(slide, Inches(6.1), Inches(6.2), Inches(6.5), Inches(0.3), "DB: SystemSetting (key + value + category)", 10, ACCENT_BLUE, False)
add_text(slide, Inches(6.1), Inches(6.5), Inches(6.5), Inches(0.3), "DB: DeviceSetting (deviceId + key + value, cascade delete)", 10, ACCENT_GREEN, False)
add_text(slide, Inches(6.1), Inches(6.8), Inches(6.5), Inches(0.3), "현재 기기 식별: POSON_DEVICE_ID 환경변수", 10, ACCENT_ORANGE, False)


# ════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "asis", "POSON_설정_항목_전체_목록.pptx")
prs.save(output_path)
print(f"PPTX saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
print(f"File size: {os.path.getsize(output_path) / 1024:.0f} KB")
