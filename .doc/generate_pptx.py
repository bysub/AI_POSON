"""ASIS 매입~판매 비즈니스 플로우 PPTX 생성"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Color Palette ───
BG_DARK = RGBColor(0x1E, 0x29, 0x3B)
BG_CARD = RGBColor(0x27, 0x34, 0x4A)
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)
ACCENT_GREEN = RGBColor(0x4A, 0xDE, 0x80)
ACCENT_ORANGE = RGBColor(0xFB, 0xBF, 0x24)
ACCENT_RED = RGBColor(0xF8, 0x71, 0x71)
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
TEXT_LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
BORDER_COLOR = RGBColor(0x33, 0x41, 0x55)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=14, color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.font.name = "Malgun Gothic"
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=12, color=TEXT_WHITE, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, text_color, is_bold, sz) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz or font_size)
        p.font.color.rgb = text_color or color
        p.font.bold = is_bold
        p.font.name = "Malgun Gothic"
        p.space_after = Pt(2)
    return txBox


def add_flow_arrow(slide, left, top, width=Inches(0.4), color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_badge(slide, left, top, text, color=ACCENT_BLUE, width=Inches(1.2)):
    shape = add_shape(slide, left, top, width, Inches(0.35), fill_color=color)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = TEXT_WHITE
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "Malgun Gothic"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


# ════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Top accent line
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_BLUE)

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
         "POSON POS 레거시 시스템", 20, TEXT_GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.6), Inches(11), Inches(1.2),
         "매입 ~ 판매 비즈니스 플로우 분석", 40, TEXT_WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
         "ASIS 분석 | VB6 레거시 소스 기반 | POS & KIOSK 모드 비교", 16, ACCENT_BLUE, False, PP_ALIGN.CENTER)

# Info cards
for i, (label, value) in enumerate([
    ("분석 대상", "300K+ LOC / 152 폼 / 26 모듈"),
    ("분석 기준", "POSON_POS, POSON_POS_SELF21"),
    ("분석일", "2026-02-19"),
]):
    x = Inches(2.5 + i * 3)
    card = add_shape(slide, x, Inches(5.2), Inches(2.6), Inches(1.1), fill_color=BG_CARD, border_color=BORDER_COLOR)
    add_text(slide, x + Inches(0.2), Inches(5.3), Inches(2.2), Inches(0.4), label, 11, TEXT_GRAY, False)
    add_text(slide, x + Inches(0.2), Inches(5.7), Inches(2.2), Inches(0.5), value, 13, TEXT_WHITE, True)


# ════════════════════════════════════════════════════════════════
# SLIDE 2: Business Cycle Overview
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_BLUE)
add_text(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.6), "전체 비즈니스 사이클", 28, TEXT_WHITE, True)

# 5 cycle boxes
cycle_items = [
    ("01", "매입", "거래처 → 입고\nInD/InT 테이블", ACCENT_BLUE),
    ("02", "상품등록", "Goods 마스터\n분류코드/가격", ACCENT_GREEN),
    ("03", "판매", "POS / KIOSK\nSaD/SaT 테이블", ACCENT_ORANGE),
    ("04", "결제", "VAN 12사\n카드/현금/포인트", ACCENT_RED),
    ("05", "정산/마감", "일계 DF 테이블\nZ리포트 출력", ACCENT_PURPLE),
]
for i, (num, title, desc, color) in enumerate(cycle_items):
    x = Inches(0.5 + i * 2.5)
    y = Inches(1.5)
    card = add_shape(slide, x, y, Inches(2.2), Inches(2.0), fill_color=BG_CARD, border_color=color)
    # Number badge
    add_badge(slide, x + Inches(0.05), y + Inches(0.1), num, color, Inches(0.5))
    add_text(slide, x + Inches(0.15), y + Inches(0.55), Inches(1.9), Inches(0.5), title, 18, TEXT_WHITE, True)
    add_text(slide, x + Inches(0.15), y + Inches(1.1), Inches(1.9), Inches(0.8), desc, 11, TEXT_GRAY, False)
    if i < 4:
        add_flow_arrow(slide, x + Inches(2.2), y + Inches(0.85), Inches(0.3), color)

# DB Tables section
add_text(slide, Inches(0.6), Inches(3.8), Inches(12), Inches(0.5), "핵심 DB 테이블 (월별 동적 생성)", 18, TEXT_WHITE, True)

table_data = [
    ("SaT_{YYYYMM}", "판매 총계", "Sale_Num(PK), 40+ 컬럼"),
    ("SaD_{YYYYMM}", "판매 상세", "Sale_Num + Sale_Seq"),
    ("InT_{YYYYMM}", "입고 총계", "In_Date 기준 집계"),
    ("InD_{YYYYMM}", "입고 상세", "Barcode + In_Date"),
    ("LastSt_{YYYYMM}", "이월재고", "BarCode(PK), Real_Sto"),
    ("DF_{YYYYMM}", "일계 정산", "Sale_Date(PK)"),
]

# Master tables
master_data = [
    ("Goods", "상품 마스터", "BarCode(PK)"),
    ("Office_Manage", "거래처", "Office_Code(PK)"),
    ("Bundle", "묶음상품", "P_Barcode + C_Barcode"),
    ("Stock_Discard", "폐기 재고", "Barcode + SD_DATE"),
]

for col, dataset, title_text, accent in [
    (0, table_data, "월별 동적 테이블", ACCENT_BLUE),
    (1, master_data, "마스터 테이블", ACCENT_GREEN),
]:
    x_base = Inches(0.5 + col * 6.3)
    add_badge(slide, x_base, Inches(4.35), title_text, accent, Inches(1.8))
    for j, (tbl, desc, pk) in enumerate(dataset):
        y_pos = Inches(4.85 + j * 0.4)
        add_text(slide, x_base + Inches(0.1), y_pos, Inches(2), Inches(0.35), tbl, 11, ACCENT_BLUE, True)
        add_text(slide, x_base + Inches(2.2), y_pos, Inches(1.5), Inches(0.35), desc, 11, TEXT_LIGHT, False)
        add_text(slide, x_base + Inches(3.8), y_pos, Inches(2.2), Inches(0.35), pk, 10, TEXT_GRAY, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 3: POS - 매입 Flow (UC-P01)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_BLUE)
add_badge(slide, Inches(0.5), Inches(0.25), "POS", ACCENT_BLUE, Inches(0.7))
add_text(slide, Inches(1.3), Inches(0.2), Inches(8), Inches(0.5), "UC-P01: 매입 (Purchase/Receiving)", 24, TEXT_WHITE, True)

# Info row
info_items = [
    ("Actor", "관리자"),
    ("핵심 폼", "Frm_PurchaseManage"),
    ("DB", "InD/InT_{YYYYMM}, LastSt_{YYYYMM}"),
]
for i, (k, v) in enumerate(info_items):
    x = Inches(0.5 + i * 4.2)
    add_text(slide, x, Inches(0.85), Inches(1), Inches(0.3), k, 10, TEXT_GRAY, True)
    add_text(slide, x + Inches(1.05), Inches(0.85), Inches(3), Inches(0.3), v, 11, ACCENT_BLUE, False)

# Main Flow steps
steps = [
    "1. 매입 관리 화면 (Frm_PurchaseManage) 진입",
    "2. 거래처 선택 (Office_Manage 테이블 조회)",
    "3. 매입일자, 매입유형 선택 (정상/반품/누락/재고조정)",
    "4. 상품 바코드 스캔 또는 검색 (Goods 테이블)",
    "5. 매입단가(Pur_Pri) × 수량(In_Count) 입력",
    "6. 입고금액 자동 계산: In_Pri = Pur_Pri × In_Count",
    "7. 박스 상품: BOX_USE='1' → OBTAIN 수량 자동 적용",
    "8. 매입 확정 → InD_{YYYYMM} INSERT (입고 상세)",
    "9. InT_{YYYYMM} UPDATE (입고 총계)",
    "10. LastSt_{YYYYMM} Real_Sto 증가 (이월재고)",
    "11. (선택) 매입 영수증 출력",
]
card = add_shape(slide, Inches(0.4), Inches(1.3), Inches(7.5), Inches(5.8), fill_color=BG_CARD, border_color=BORDER_COLOR)
for j, step in enumerate(steps):
    color = ACCENT_GREEN if "INSERT" in step or "UPDATE" in step or "증가" in step else TEXT_LIGHT
    add_text(slide, Inches(0.7), Inches(1.5 + j * 0.48), Inches(7), Inches(0.45), step, 13, color, False)

# Side: 이월재고 공식
add_text(slide, Inches(8.3), Inches(1.3), Inches(4.5), Inches(0.4), "이월재고 계산 공식", 16, TEXT_WHITE, True)
formula_card = add_shape(slide, Inches(8.2), Inches(1.8), Inches(4.7), Inches(3.5), fill_color=RGBColor(0x1A, 0x24, 0x35), border_color=ACCENT_BLUE)
formula_lines = [
    ("Real_Sto =", ACCENT_BLUE, True, 13),
    ("  전월이월", TEXT_LIGHT, False, 12),
    ("  - 판매수량 (SaD_Count)", ACCENT_ORANGE, False, 12),
    ("  - 반품수량 (SdD_Count)", ACCENT_ORANGE, False, 12),
    ("  + 입고수량 (InD_Count)", ACCENT_GREEN, False, 12),
    ("  + 재고조정 (StSet_Count)", ACCENT_GREEN, False, 12),
    ("  - 폐기수량 (SD_COUNT)", ACCENT_RED, False, 12),
    ("  - 박스판매 (BP_SaD_Count)", ACCENT_ORANGE, False, 12),
    ("  + 박스입고 (BP_InD_Count)", ACCENT_GREEN, False, 12),
]
add_multiline(slide, Inches(8.5), Inches(2.0), Inches(4.2), Inches(3.2), formula_lines)

# Alternative Flow
add_text(slide, Inches(8.3), Inches(5.5), Inches(4.5), Inches(0.4), "Alternative Flow", 14, TEXT_WHITE, True)
alt_lines = [
    ("4a  엑셀 대량 매입", TEXT_LIGHT, False, 11),
    ("     → Frm_PurchaseExcelCall 일괄 입력", TEXT_GRAY, False, 10),
    ("5a  신규 상품", TEXT_LIGHT, False, 11),
    ("     → Frm_Purchase_New_Product 등록 후 매입", TEXT_GRAY, False, 10),
    ("8a  매입 취소", TEXT_LIGHT, False, 11),
    ("     → 재고 감소, InD 삭제", ACCENT_RED, False, 10),
]
add_multiline(slide, Inches(8.5), Inches(5.9), Inches(4.2), Inches(1.5), alt_lines)


# ════════════════════════════════════════════════════════════════
# SLIDE 4: POS - 판매 Flow (UC-P03)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_ORANGE)
add_badge(slide, Inches(0.5), Inches(0.25), "POS", ACCENT_BLUE, Inches(0.7))
add_text(slide, Inches(1.3), Inches(0.2), Inches(10), Inches(0.5), "UC-P03: POS 판매 Main Flow", 24, TEXT_WHITE, True)

info_items2 = [
    ("Actor", "판매원 (캐셔)"),
    ("핵심 폼", "Frm_SaleMain (~16,000줄)"),
    ("DB", "SaT/SaD_{YYYYMM}, LastSt_{YYYYMM}"),
]
for i, (k, v) in enumerate(info_items2):
    x = Inches(0.5 + i * 4.2)
    add_text(slide, x, Inches(0.85), Inches(1), Inches(0.3), k, 10, TEXT_GRAY, True)
    add_text(slide, x + Inches(1.05), Inches(0.85), Inches(3), Inches(0.3), v, 11, ACCENT_ORANGE, False)

# Left column: Steps 1-10
left_steps = [
    ("1. POS 판매 화면 진입", TEXT_LIGHT),
    ("2. 영업일 결정 (새벽 4시 기준)", TEXT_LIGHT),
    ("3. 월별 동적 테이블 확인/생성", TEXT_GRAY),
    ("4. HW 초기화 (CDP, 스캐너, MSR)", TEXT_GRAY),
    ("5. 상품 바코드 스캔 / 핫키 / 검색", ACCENT_BLUE),
    ("6. Goods 테이블 상품 정보 조회", TEXT_LIGHT),
    ("7. 판매 그리드에 상품 추가", TEXT_LIGHT),
    ("8. 박스 상품 OBTAIN 수량 적용", TEXT_GRAY),
    ("9. (반복) 추가 상품 선택", TEXT_LIGHT),
    ("10. 회원 할인 적용 (MEM_SalePrice)", ACCENT_GREEN),
]
card1 = add_shape(slide, Inches(0.3), Inches(1.3), Inches(6.2), Inches(5.8), fill_color=BG_CARD, border_color=BORDER_COLOR)
for j, (step, color) in enumerate(left_steps):
    add_text(slide, Inches(0.6), Inches(1.5 + j * 0.48), Inches(5.8), Inches(0.45), step, 13, color, False)

# Right column: Steps 11-19
right_steps = [
    ("11. 금액/비율 할인 적용", TEXT_LIGHT),
    ("12. 총금액 계산: 소계-할인+봉사료", ACCENT_ORANGE),
    ("13. 결제 수단 선택", ACCENT_BLUE),
    ("     현금 / 카드 / 포인트 / 상품권", TEXT_GRAY),
    ("14. 결제 처리 → UC-P04", ACCENT_BLUE),
    ("15. SaT_{YYYYMM} INSERT (판매총계)", ACCENT_GREEN),
    ("16. SaD_{YYYYMM} INSERT (판매상세)", ACCENT_GREEN),
    ("17. LastSt Real_Sto 차감 (재고)", ACCENT_RED),
    ("18. 영수증 출력 (Mdl_Printer)", TEXT_LIGHT),
    ("19. 현금 시 서랍 자동 개방", TEXT_GRAY),
]
card2 = add_shape(slide, Inches(6.7), Inches(1.3), Inches(6.2), Inches(5.8), fill_color=BG_CARD, border_color=BORDER_COLOR)
for j, (step, color) in enumerate(right_steps):
    add_text(slide, Inches(7.0), Inches(1.5 + j * 0.48), Inches(5.8), Inches(0.45), step, 13, color, False)

# Sale_Num format
add_text(slide, Inches(7.0), Inches(6.3), Inches(5.5), Inches(0.3), "전표번호: POS번호(2) + 년월일(8) + 시분초(6) + 순번(5)", 10, TEXT_GRAY, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 5: POS - 결제 Flow (UC-P04)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_RED)
add_badge(slide, Inches(0.5), Inches(0.25), "POS", ACCENT_BLUE, Inches(0.7))
add_text(slide, Inches(1.3), Inches(0.2), Inches(10), Inches(0.5), "UC-P04: 결제 처리 (VAN 12사 연동)", 24, TEXT_WHITE, True)

# Card payment flow
add_text(slide, Inches(0.5), Inches(1.0), Inches(6), Inches(0.4), "카드 결제 Main Flow", 18, TEXT_WHITE, True)
card_steps = [
    ("1. VAN사 확인 (VAN_SELECT 설정값)", TEXT_LIGHT),
    ("2. VAN별 결제 폼 로드", TEXT_LIGHT),
    ("3. 카드 리딩 (MSR 스와이프 / IC칩 / 키인)", ACCENT_BLUE),
    ("4. 할부 개월 선택 (00=일시불 ~ 12)", TEXT_LIGHT),
    ("5. 승인 요청 전문 생성 (TID/금액/카드정보)", ACCENT_ORANGE),
    ("6. 서명패드 전자서명 획득 (5만원 이상)", TEXT_GRAY),
    ("7. VAN DLL 호출 → VAN 서버 전송", ACCENT_RED),
    ("8. 카드사 승인/거절 응답 수신", ACCENT_GREEN),
    ("9. CARD_TRAN 테이블 INSERT", ACCENT_GREEN),
    ("10. SaT.Card_Pri UPDATE + 영수증 출력", TEXT_LIGHT),
]
card_bg = add_shape(slide, Inches(0.3), Inches(1.5), Inches(6.5), Inches(5.5), fill_color=BG_CARD, border_color=BORDER_COLOR)
for j, (step, color) in enumerate(card_steps):
    add_text(slide, Inches(0.6), Inches(1.7 + j * 0.47), Inches(6), Inches(0.42), step, 12, color, False)

# VAN list
add_text(slide, Inches(7.2), Inches(1.0), Inches(5.5), Inches(0.4), "지원 VAN 12사", 18, TEXT_WHITE, True)
van_list = [
    ("NICE", "NicePosV205.dll", "SSL"),
    ("KSNET", "KSNet_ADSL.dll", "TCP/IP"),
    ("KICC", "Kicc.dll + KiccDSC.ocx", "OCX"),
    ("KIS", "KisCatSSL.dll", "SSL+BC"),
    ("SMARTRO", "SmartroSign.dll", "NFC/T-money"),
    ("FDIK", "fdikpos43.dll", "현금영수증"),
    ("JTNET", "NCPOS.dll", "서명패드"),
    ("KCP", "KCPOCX.ocx", "OCX"),
    ("KOCES", "API 방식", "-"),
    ("KOVAN", "kovan_signpad.ocx", "OCX"),
    ("SPC", "SPCNSecuCAT.ocx", "OCX"),
    ("StarVAN", "StarVAN.ocx", "OCX"),
]
van_bg = add_shape(slide, Inches(7.0), Inches(1.5), Inches(5.9), Inches(5.5), fill_color=BG_CARD, border_color=BORDER_COLOR)
# Header
add_text(slide, Inches(7.3), Inches(1.6), Inches(1.3), Inches(0.35), "VAN사", 10, TEXT_GRAY, True)
add_text(slide, Inches(8.6), Inches(1.6), Inches(2.8), Inches(0.35), "DLL/OCX", 10, TEXT_GRAY, True)
add_text(slide, Inches(11.5), Inches(1.6), Inches(1.2), Inches(0.35), "방식", 10, TEXT_GRAY, True)
for j, (van, dll, method) in enumerate(van_list):
    y = Inches(2.0 + j * 0.38)
    color = ACCENT_BLUE if j < 6 else TEXT_LIGHT
    add_text(slide, Inches(7.3), y, Inches(1.3), Inches(0.35), van, 11, color, True)
    add_text(slide, Inches(8.6), y, Inches(2.8), Inches(0.35), dll, 10, TEXT_GRAY, False)
    add_text(slide, Inches(11.5), y, Inches(1.2), Inches(0.35), method, 10, TEXT_GRAY, False)

# Cash payment
add_text(slide, Inches(7.3), Inches(6.0), Inches(5), Inches(0.3), "현금 결제: 투입 → 거스름돈 → 현금영수증(선택) → 서랍개방", 11, ACCENT_GREEN, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 6: KIOSK Flow (UC-K01) - Part 1
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_GREEN)
add_badge(slide, Inches(0.5), Inches(0.25), "KIOSK", ACCENT_GREEN, Inches(0.9))
add_text(slide, Inches(1.5), Inches(0.2), Inches(10), Inches(0.5), "UC-K01: 키오스크 주문 플로우 (1/2)", 24, TEXT_WHITE, True)

info_kiosk = [
    ("Actor", "고객 (셀프)"),
    ("핵심 폼", "Frm_SelfKiosk (~4,700줄)"),
    ("해상도", "1080 x 1920 (세로)"),
]
for i, (k, v) in enumerate(info_kiosk):
    x = Inches(0.5 + i * 4.2)
    add_text(slide, x, Inches(0.85), Inches(1), Inches(0.3), k, 10, TEXT_GRAY, True)
    add_text(slide, x + Inches(1.05), Inches(0.85), Inches(3), Inches(0.3), v, 11, ACCENT_GREEN, False)

# Kiosk flow - visual steps
kiosk_phases = [
    ("1단계", "대기화면", "imgStartPage\n'화면을 터치하세요'\n터치 시 메인화면", ACCENT_BLUE),
    ("2단계", "카테고리", "탭 4개 + 좌우 스크롤\nlabSelfTab(0~3)\nL_Code 기준 조회", ACCENT_GREEN),
    ("3단계", "상품선택", "12개/페이지 그리드\nImage1(0~11) 터치\n품절: imgSoldOut", ACCENT_ORANGE),
    ("4단계", "옵션선택", "frm_SelfGoods 모달\n사이즈/토핑/샷\nOrderKiosk.sOPBarcods", ACCENT_PURPLE),
]

for i, (phase, title, desc, color) in enumerate(kiosk_phases):
    x = Inches(0.3 + i * 3.2)
    y = Inches(1.4)
    card = add_shape(slide, x, y, Inches(2.9), Inches(2.5), fill_color=BG_CARD, border_color=color)
    add_badge(slide, x + Inches(0.1), y + Inches(0.1), phase, color, Inches(0.9))
    add_text(slide, x + Inches(0.15), y + Inches(0.55), Inches(2.6), Inches(0.4), title, 18, TEXT_WHITE, True)
    # desc lines
    for j, line in enumerate(desc.split("\n")):
        add_text(slide, x + Inches(0.15), y + Inches(1.05 + j * 0.4), Inches(2.6), Inches(0.35), line, 11, TEXT_GRAY, False)
    if i < 3:
        add_flow_arrow(slide, x + Inches(2.9), y + Inches(1.1), Inches(0.3), color)

# Cart section
kiosk_phases2 = [
    ("5단계", "장바구니", "fg21List (VSFlexGrid)\n20컬럼: 상품명/수량/합계\n[-] [+] [x] 수량 조정", ACCENT_BLUE),
    ("6단계", "결제선택", "'결제하기' 버튼 터치\nFrm_SelfOrderList\n카드 또는 현금 선택", ACCENT_RED),
    ("7단계", "결제처리", "카드: frm_SelfCard\n현금: Frm_SelfCash\nVAN DLL 호출 (POS 동일)", ACCENT_RED),
    ("8단계", "DB저장", "SaT/SaD INSERT\nLastSt 재고 차감\nCARD_TRAN 거래 로그", ACCENT_GREEN),
]
for i, (phase, title, desc, color) in enumerate(kiosk_phases2):
    x = Inches(0.3 + i * 3.2)
    y = Inches(4.2)
    card = add_shape(slide, x, y, Inches(2.9), Inches(2.5), fill_color=BG_CARD, border_color=color)
    add_badge(slide, x + Inches(0.1), y + Inches(0.1), phase, color, Inches(0.9))
    add_text(slide, x + Inches(0.15), y + Inches(0.55), Inches(2.6), Inches(0.4), title, 18, TEXT_WHITE, True)
    for j, line in enumerate(desc.split("\n")):
        add_text(slide, x + Inches(0.15), y + Inches(1.05 + j * 0.4), Inches(2.6), Inches(0.35), line, 11, TEXT_GRAY, False)
    if i < 3:
        add_flow_arrow(slide, x + Inches(2.9), y + Inches(1.1), Inches(0.3), color)


# ════════════════════════════════════════════════════════════════
# SLIDE 7: KIOSK Flow - Part 2 (Complete + Kitchen)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_GREEN)
add_badge(slide, Inches(0.5), Inches(0.25), "KIOSK", ACCENT_GREEN, Inches(0.9))
add_text(slide, Inches(1.5), Inches(0.2), Inches(10), Inches(0.5), "UC-K01: 키오스크 주문 플로우 (2/2) - 완료 & 복귀", 24, TEXT_WHITE, True)

# Final steps
final_steps = [
    ("9단계", "영수증 출력", "Mdl_Printer\nESC/POS 프린터\n상품목록+합계+승인번호", ACCENT_ORANGE),
    ("10단계", "주방 전송", "네트워크/DB 폴링\n주문 데이터 전송\n주방 디스플레이 표시", ACCENT_PURPLE),
    ("11단계", "완료 화면", "frm_SelfPrint\n완료 GIF 애니메이션\n음성: '승인완료되었습니다'", ACCENT_BLUE),
    ("12단계", "자동 복귀", "30초 카운트다운\nFrm_SelfEnd → initmain()\n장바구니 초기화 → 대기화면", ACCENT_GREEN),
]
for i, (phase, title, desc, color) in enumerate(final_steps):
    x = Inches(0.3 + i * 3.2)
    y = Inches(1.2)
    card = add_shape(slide, x, y, Inches(2.9), Inches(2.5), fill_color=BG_CARD, border_color=color)
    add_badge(slide, x + Inches(0.1), y + Inches(0.1), phase, color, Inches(0.9))
    add_text(slide, x + Inches(0.15), y + Inches(0.55), Inches(2.6), Inches(0.4), title, 18, TEXT_WHITE, True)
    for j, line in enumerate(desc.split("\n")):
        add_text(slide, x + Inches(0.15), y + Inches(1.05 + j * 0.4), Inches(2.6), Inches(0.35), line, 11, TEXT_GRAY, False)
    if i < 3:
        add_flow_arrow(slide, x + Inches(2.9), y + Inches(1.1), Inches(0.3), color)

# Alternative Flow
add_text(slide, Inches(0.5), Inches(4.0), Inches(12), Inches(0.4), "Alternative Flow", 18, TEXT_WHITE, True)

alt_card = add_shape(slide, Inches(0.3), Inches(4.5), Inches(12.5), Inches(2.5), fill_color=BG_CARD, border_color=BORDER_COLOR)
alt_flows = [
    ("4a", "상품 품절", "imgSoldOut 오버레이 표시, 선택 불가", ACCENT_RED),
    ("5a", "옵션 없는 상품", "옵션 모달 건너뛰고 바로 장바구니 추가", TEXT_LIGHT),
    ("10a", "바코드 스캔", "txt_Input_KeyDown → addBarcode() 호출", TEXT_LIGHT),
    ("11a", "VAN 타임아웃", "자동 재시도, 실패 시 에러 메시지", ACCENT_ORANGE),
    ("11b", "카드 잔액 부족", "'잔액 부족' 메시지 → 결제 수단 재선택", ACCENT_ORANGE),
    ("12a", "네트워크 장애 (현금)", "Access MDB 오프라인 모드, 동기화 대기열", ACCENT_RED),
    ("12b", "네트워크 장애 (카드)", "'카드 결제 불가' 메시지, 현금 결제 안내", ACCENT_RED),
    ("1a", "관리자 진입", "더블클릭 → 비밀번호 → pic관리자 패널", ACCENT_PURPLE),
]
for j, (code, cond, action, color) in enumerate(alt_flows):
    y = Inches(4.65 + j * 0.28)
    add_text(slide, Inches(0.6), y, Inches(0.6), Inches(0.25), code, 10, color, True)
    add_text(slide, Inches(1.2), y, Inches(2.2), Inches(0.25), cond, 10, TEXT_LIGHT, False)
    add_text(slide, Inches(3.5), y, Inches(9), Inches(0.25), action, 10, TEXT_GRAY, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 8: POS vs KIOSK Comparison
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_PURPLE)
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "POS vs KIOSK 비교", 28, TEXT_WHITE, True)

comparison = [
    ("항목", "POS (Self_YN=0)", "KIOSK (Self_YN=1)"),
    ("Actor", "판매원 (캐셔)", "고객 (셀프)"),
    ("메인 폼", "Frm_SaleMain (~16,000줄)", "Frm_SelfKiosk (~4,700줄)"),
    ("상품 선택", "바코드 스캔 + 핫키 + 검색", "터치 그리드 (12개/페이지)"),
    ("카테고리", "대/중/소 분류 트리", "탭 4개 + 좌우 스크롤"),
    ("결제 수단", "현금+카드+포인트+상품권", "카드 + 현금"),
    ("해상도", "1024 x 768 (가로)", "1080 x 1920 (세로)"),
    ("매입/정산", "가능 (전용 폼)", "불가 (POS에서만)"),
    ("대기화면", "없음 (항상 판매)", "imgStartPage (터치 대기)"),
    ("자동 복귀", "없음", "30초 카운트다운"),
    ("주방 전송", "수동/자동", "자동"),
    ("관리자", "메뉴 기반 전체 접근", "더블클릭 → 제한된 관리"),
    ("하드웨어", "CDP+스캐너+MSR+저울+서랍", "터치모니터+카드리더"),
]

# Table header
header_y = Inches(1.2)
add_shape(slide, Inches(0.5), header_y, Inches(2.8), Inches(0.5), fill_color=RGBColor(0x33, 0x41, 0x55))
add_shape(slide, Inches(3.4), header_y, Inches(4.5), Inches(0.5), fill_color=ACCENT_BLUE)
add_shape(slide, Inches(8.0), header_y, Inches(4.5), Inches(0.5), fill_color=ACCENT_GREEN)

for j, (item, pos_val, kiosk_val) in enumerate(comparison):
    y = header_y + Inches(j * 0.45)
    if j == 0:
        add_text(slide, Inches(0.7), y + Inches(0.07), Inches(2.5), Inches(0.4), item, 12, TEXT_WHITE, True)
        add_text(slide, Inches(3.6), y + Inches(0.07), Inches(4.2), Inches(0.4), pos_val, 12, TEXT_WHITE, True)
        add_text(slide, Inches(8.2), y + Inches(0.07), Inches(4.2), Inches(0.4), kiosk_val, 12, TEXT_WHITE, True)
    else:
        if j % 2 == 0:
            add_shape(slide, Inches(0.5), y, Inches(12), Inches(0.45), fill_color=RGBColor(0x1A, 0x24, 0x35))
        add_text(slide, Inches(0.7), y + Inches(0.07), Inches(2.5), Inches(0.4), item, 11, TEXT_GRAY, True)
        add_text(slide, Inches(3.6), y + Inches(0.07), Inches(4.2), Inches(0.4), pos_val, 11, ACCENT_BLUE, False)
        add_text(slide, Inches(8.2), y + Inches(0.07), Inches(4.2), Inches(0.4), kiosk_val, 11, ACCENT_GREEN, False)


# ════════════════════════════════════════════════════════════════
# SLIDE 9: Migration Mapping
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_PURPLE)
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "마이그레이션 매핑: VB6 → Electron + Vue 3", 28, TEXT_WHITE, True)

mappings = [
    ("Frm_SaleMain", "SalesRegisterView.vue", "POS 판매"),
    ("Frm_SelfKiosk", "MenuView.vue + CartView.vue", "키오스크"),
    ("Frm_PurchaseManage", "PurchaseRegisterView.vue", "매입 관리"),
    ("Frm_ProductUpdate", "ProductsView.vue", "상품 관리"),
    ("Frm_SettlementManage", "DashboardView.vue (정산 탭)", "정산/마감"),
    ("Frm_SaleCard_VCAT_*", "PaymentService (Strategy)", "12 VAN → 4 VAN"),
    ("InD/InT/SaD/SaT (월별)", "Purchase/Order (단일 테이블)", "Prisma 모델"),
    ("LastSt (이월재고)", "PurchaseProduct.stock", "재고 필드"),
    ("Goods (마스터)", "Product + Category (M:N)", "정규화"),
    ("Office_Manage", "Supplier", "거래처"),
    ("Mdl_Printer", "Electron IPC → ESC/POS", "프린터 서비스"),
    ("pos_config.ini / POS_Set", "SystemSetting + DeviceSetting", "2계층 설정"),
]

# Header
my = Inches(1.2)
add_shape(slide, Inches(0.5), my, Inches(4), Inches(0.5), fill_color=ACCENT_RED)
add_shape(slide, Inches(4.6), my, Inches(0.3), Inches(0.5), fill_color=BG_DARK)
add_shape(slide, Inches(5.0), my, Inches(5), Inches(0.5), fill_color=ACCENT_GREEN)
add_shape(slide, Inches(10.1), my, Inches(0.3), Inches(0.5), fill_color=BG_DARK)
add_shape(slide, Inches(10.5), my, Inches(2.3), Inches(0.5), fill_color=ACCENT_PURPLE)

add_text(slide, Inches(0.7), my + Inches(0.07), Inches(3.8), Inches(0.4), "VB6 ASIS", 13, TEXT_WHITE, True)
add_text(slide, Inches(5.2), my + Inches(0.07), Inches(4.8), Inches(0.4), "Vue 3 TOBE", 13, TEXT_WHITE, True)
add_text(slide, Inches(10.7), my + Inches(0.07), Inches(2), Inches(0.4), "비고", 13, TEXT_WHITE, True)

for j, (asis, tobe, note) in enumerate(mappings):
    y = my + Inches(0.55 + j * 0.45)
    if j % 2 == 0:
        add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(0.45), fill_color=RGBColor(0x1A, 0x24, 0x35))
    add_text(slide, Inches(0.7), y + Inches(0.07), Inches(3.8), Inches(0.35), asis, 12, ACCENT_RED, False)
    add_text(slide, Inches(4.7), y + Inches(0.12), Inches(0.2), Inches(0.25), "→", 11, TEXT_GRAY, True)
    add_text(slide, Inches(5.2), y + Inches(0.07), Inches(4.8), Inches(0.35), tobe, 12, ACCENT_GREEN, False)
    add_text(slide, Inches(10.7), y + Inches(0.07), Inches(2), Inches(0.35), note, 10, TEXT_GRAY, False)


# ════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "asis", "POSON_매입_판매_플로우_분석.pptx")
prs.save(output_path)
print(f"PPTX saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
